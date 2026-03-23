import asyncio
import json
import io
import os
import random
import secrets
import sqlite3
import string
import textwrap
import time
import uuid
import zipfile
from base64 import b64decode
from io import BytesIO
from hashlib import pbkdf2_hmac, sha256
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any
from urllib.error import HTTPError, URLError
from urllib.request import Request, urlopen

from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, Header, HTTPException, Query, Response, UploadFile, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, RedirectResponse, Response
from urllib.parse import urlencode, urlparse, parse_qs
from pydantic import BaseModel
from app.assessment_reporting import (
    AssessmentResponse,
    build_individual_profiles,
    calculate_session_metrics,
    fetch_session_responses,
    init_assessment_schema,
    question_performance_breakdown,
    save_question_performance_bar_chart,
    save_student_summary_table,
    store_assessment_response,
)

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.units import cm
    from reportlab.pdfbase.pdfmetrics import stringWidth
    from reportlab.pdfgen import canvas
except Exception:  # pragma: no cover
    A4 = None
    ImageReader = None
    cm = None
    stringWidth = None
    canvas = None

try:
    from openai import OpenAI
except Exception:  # pragma: no cover
    OpenAI = None


BASE_DIR = Path(__file__).resolve().parents[1]
# Always read the backend-local .env file, regardless of the process working directory.
load_dotenv(BASE_DIR / ".env")
DB_PATH = BASE_DIR / "data.sqlite3"
UPLOADS_DIR = BASE_DIR / "uploads"

# Timeout for AI quiz generation calls (seconds). Long enough for slow providers.
AI_QUIZ_GENERATION_TIMEOUT_SECONDS = 45.0

# OAuth configuration (set in backend/.env — never commit actual secrets)
GITHUB_CLIENT_ID = os.getenv("GITHUB_CLIENT_ID", "").strip()
GITHUB_CLIENT_SECRET = os.getenv("GITHUB_CLIENT_SECRET", "").strip()
FRONTEND_URL = os.getenv("FRONTEND_URL", "http://localhost:5173").strip().rstrip("/")


def parse_bearer_token(authorization: str | None) -> str | None:
    if not authorization:
        return None
    parts = authorization.strip().split(" ", 1)
    if len(parts) != 2 or parts[0].lower() != "bearer":
        return None
    token = parts[1].strip()
    return token or None


def hash_password(password: str, salt_hex: str | None = None) -> tuple[str, str]:
    salt = bytes.fromhex(salt_hex) if salt_hex else secrets.token_bytes(16)
    password_hash = pbkdf2_hmac("sha256", password.encode("utf-8"), salt, 120_000)
    return salt.hex(), password_hash.hex()


def create_auth_token() -> str:
    return secrets.token_urlsafe(32)


def upsert_oauth_user(email: str, display_name: str) -> tuple[dict[str, Any], str]:
    """Find or create a user by email (for OAuth sign-in), return (user_dict, token)."""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute("SELECT id, email, display_name, role FROM users WHERE email = ?", (email,))
    row = cursor.fetchone()
    if row:
        user_id = row["id"]
        user = dict(row)
    else:
        cursor.execute(
            "INSERT INTO users(email, display_name, role, password_salt, password_hash, created_at) VALUES (?, ?, ?, ?, ?, ?)",
            (email, display_name, "teacher", "", "", now_iso()),
        )
        user_id = cursor.lastrowid
        cursor.execute("SELECT id, email, display_name, role FROM users WHERE id = ?", (user_id,))
        user = dict(cursor.fetchone())
    token = create_auth_token()
    cursor.execute(
        "INSERT INTO auth_tokens(token, user_id, created_at) VALUES (?, ?, ?)",
        (token, user_id, now_iso()),
    )
    conn.commit()
    conn.close()
    return user, token


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def init_db() -> None:
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS sessions (
            code TEXT PRIMARY KEY,
            created_at TEXT NOT NULL,
            teacher_name TEXT NOT NULL,
            owner_user_id INTEGER,
            active INTEGER NOT NULL DEFAULT 1
        )
        """
    )
    cursor.execute("PRAGMA table_info(sessions)")
    session_columns = {row[1] for row in cursor.fetchall()}
    if "owner_user_id" not in session_columns:
        cursor.execute("ALTER TABLE sessions ADD COLUMN owner_user_id INTEGER")
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS events (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            session_code TEXT NOT NULL,
            event_type TEXT NOT NULL,
            payload TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(session_code) REFERENCES sessions(code)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email TEXT NOT NULL UNIQUE,
            display_name TEXT NOT NULL,
            role TEXT NOT NULL,
            password_salt TEXT NOT NULL,
            password_hash TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS auth_tokens (
            token TEXT PRIMARY KEY,
            user_id INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(user_id) REFERENCES users(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS presentations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            session_code TEXT,
            original_name TEXT NOT NULL,
            stored_name TEXT NOT NULL,
            mime_type TEXT,
            size_bytes INTEGER NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(user_id) REFERENCES users(id)
        )
        """
    )
    cursor.execute("PRAGMA table_info(presentations)")
    presentation_columns = {row[1] for row in cursor.fetchall()}
    if "session_code" not in presentation_columns:
        cursor.execute("ALTER TABLE presentations ADD COLUMN session_code TEXT")
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS saved_quizzes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            session_code TEXT,
            question TEXT NOT NULL,
            options_json TEXT NOT NULL,
            correct_option_id TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(user_id) REFERENCES users(id)
        )
        """
    )
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS ai_cache (
            cache_key TEXT PRIMARY KEY,
            response TEXT NOT NULL,
            created_at TEXT NOT NULL,
            expires_at TEXT NOT NULL,
            hit_count INTEGER NOT NULL DEFAULT 0
        )
        """
    )
    conn.commit()
    conn.close()

    # Reporting schema is maintained in a dedicated module.
    init_assessment_schema(DB_PATH)


def ai_cache_key(prefix: str, *parts: str) -> str:
    """Return a SHA-256 hex digest over the concatenation of prefix and all parts."""
    payload = prefix + "|" + "|".join(parts)
    return sha256(payload.encode("utf-8")).hexdigest()


def ai_cache_get(cache_key: str) -> str | None:
    """Return cached response string, or None on miss / expiry."""
    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.cursor()
        cursor.execute(
            "SELECT response FROM ai_cache WHERE cache_key = ? AND expires_at > ?",
            (cache_key, now_iso()),
        )
        row = cursor.fetchone()
        if row:
            cursor.execute(
                "UPDATE ai_cache SET hit_count = hit_count + 1 WHERE cache_key = ?",
                (cache_key,),
            )
            conn.commit()
            return row[0]
        return None
    finally:
        conn.close()


def ai_cache_set(cache_key: str, response: str, ttl_seconds: int = 3600) -> None:
    """Store response in cache with a TTL. Overwrites any existing entry."""
    from datetime import timedelta
    conn = sqlite3.connect(DB_PATH)
    try:
        cursor = conn.cursor()
        now = datetime.now(timezone.utc)
        expires = (now + timedelta(seconds=ttl_seconds)).isoformat()
        cursor.execute(
            """
            INSERT INTO ai_cache(cache_key, response, created_at, expires_at, hit_count)
            VALUES (?, ?, ?, ?, 0)
            ON CONFLICT(cache_key) DO UPDATE SET
                response = excluded.response,
                created_at = excluded.created_at,
                expires_at = excluded.expires_at,
                hit_count = 0
            """,
            (cache_key, response, now.isoformat(), expires),
        )
        conn.commit()
    finally:
        conn.close()


def insert_session(code: str, teacher_name: str, owner_user_id: int | None = None) -> None:
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        "INSERT INTO sessions(code, created_at, teacher_name, owner_user_id, active) VALUES (?, ?, ?, ?, 1)",
        (code, now_iso(), teacher_name, owner_user_id),
    )
    conn.commit()
    conn.close()


def insert_event(session_code: str, event_type: str, payload: dict[str, Any]) -> None:
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        "INSERT INTO events(session_code, event_type, payload, created_at) VALUES (?, ?, ?, ?)",
        (session_code, event_type, json.dumps(payload), now_iso()),
    )
    conn.commit()
    conn.close()


def get_user_by_token(token: str | None) -> dict[str, Any] | None:
    if not token:
        return None
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT u.id, u.email, u.display_name, u.role
        FROM auth_tokens t
        JOIN users u ON u.id = t.user_id
        WHERE t.token = ?
        """,
        (token,),
    )
    row = cursor.fetchone()
    conn.close()
    return dict(row) if row else None


def require_user(authorization: str | None) -> dict[str, Any]:
    token = parse_bearer_token(authorization)
    user = get_user_by_token(token)
    if not user:
        raise HTTPException(status_code=401, detail="Unauthorized")
    return user


def set_session_active(code: str, active: bool) -> None:
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("UPDATE sessions SET active = ? WHERE code = ?", (1 if active else 0, code))
    conn.commit()
    conn.close()


def session_exists(code: str) -> bool:
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    # Code must be globally unique because sessions.code is a primary key.
    cursor.execute("SELECT code FROM sessions WHERE code = ?", (code,))
    row = cursor.fetchone()
    conn.close()
    return row is not None


class SessionCreateRequest(BaseModel):
    teacher_name: str


class SessionCreateResponse(BaseModel):
    code: str


class AuthRegisterRequest(BaseModel):
    email: str
    display_name: str
    password: str
    role: str = "teacher"


class AuthLoginRequest(BaseModel):
    email: str
    password: str


class UserPublic(BaseModel):
    id: int
    email: str
    display_name: str
    role: str


class AuthResponse(BaseModel):
    token: str
    user: UserPublic


class PresentationItem(BaseModel):
    id: int
    session_code: str | None
    original_name: str
    mime_type: str
    size_bytes: int
    created_at: str
    download_url: str


class QuizOption(BaseModel):
    id: str
    text: str


class QuizPayload(BaseModel):
    question: str
    options: list[QuizOption]
    correct_option_id: str


class SavedQuizCreateRequest(BaseModel):
    session_code: str | None = None
    question: str
    options: list[QuizOption]
    correct_option_id: str


class SavedQuizItem(BaseModel):
    id: int
    session_code: str | None
    question: str
    options: list[QuizOption]
    correct_option_id: str
    created_at: str

class SavedQuizListItem(BaseModel):
    id: int
    session_code: str | None
    question: str
    options: list[QuizOption]
    correct_option_id: str | None
    answer_revealed: bool
    is_live: bool
    created_at: str


@dataclass
class ClientState:
    client_id: str
    websocket: WebSocket
    role: str
    name: str
    joined_at: float = field(default_factory=time.time)
    last_break_vote_at: float = 0.0
    last_confusion_vote_at: float | None = None
    confusion_signals_sent: int = 0
    break_votes_cast: int = 0
    quiz_answers_submitted: int = 0
    quiz_correct_answers: int = 0
    last_active_at: float = field(default_factory=time.time)


@dataclass
class AnonymousQuestion:
    id: str
    text: str
    created_at: str
    resolved: bool = False
    resolved_at: str | None = None


FOCUS_PERIOD_SECONDS = 1800  # 30 minutes of teaching before break-vote button unlocks


@dataclass
class RuntimeSession:
    code: str
    teacher_name: str
    created_at_epoch: float = field(default_factory=time.time)
    clients: dict[str, ClientState] = field(default_factory=dict)
    break_votes: set[str] = field(default_factory=set)
    break_active_until: float | None = None
    # Epoch when break-vote button unlocks for students; 0.0 = already unlocked
    focus_period_ends_at: float = 0.0
    notes: str = ""
    current_quiz: QuizPayload | None = None
    current_quiz_saved_id: int | None = None
    quiz_answers: dict[str, str] = field(default_factory=dict)
    quiz_hidden: bool = False
    quiz_cover_mode: bool = True
    quiz_voting_closed: bool = False
    quiz_answer_revealed: bool = False
    anonymous_questions: list[AnonymousQuestion] = field(default_factory=list)
    active: bool = True
    ended_at_epoch: float | None = None
    final_report: dict[str, Any] | None = None
    final_report_insights: dict[str, Any] | None = None
    engagement_timeline: list[dict[str, Any]] = field(default_factory=list)

    def __post_init__(self) -> None:
        # Set focus period relative to session creation if not already set
        if self.focus_period_ends_at == 0.0:
            self.focus_period_ends_at = self.created_at_epoch + FOCUS_PERIOD_SECONDS

    @property
    def student_count(self) -> int:
        return sum(1 for c in self.clients.values() if c.role == "student")


CONFUSION_DECAY_SECONDS = 90.0
CONFUSION_ACTIVE_THRESHOLD = 0.2


def current_student_confusion_level(student: ClientState, now_epoch: float) -> float:
    if student.role != "student" or student.last_confusion_vote_at is None:
        return 0.0

    elapsed = max(0.0, now_epoch - student.last_confusion_vote_at)
    return max(0.0, 1.0 - (elapsed / CONFUSION_DECAY_SECONDS))


def confusion_snapshot(session: RuntimeSession, now_epoch: float | None = None) -> dict[str, Any]:
    now_epoch = now_epoch if now_epoch is not None else time.time()
    total_level = 0.0
    active_students = 0
    student_count = 0

    for client in session.clients.values():
        if client.role != "student":
            continue
        student_count += 1
        level = current_student_confusion_level(client, now_epoch)
        total_level += level
        if level >= CONFUSION_ACTIVE_THRESHOLD:
            active_students += 1

    average_level = (total_level / student_count) if student_count else 0.0
    level_percent = int(round(average_level * 100))
    return {
        "average_level": average_level,
        "level_percent": max(0, min(level_percent, 100)),
        "active_students": active_students,
        "student_count": student_count,
    }


def metrics_payload(session: RuntimeSession) -> dict[str, Any]:
    confusion = confusion_snapshot(session)
    student_count = confusion["student_count"]
    ratio = (len(session.break_votes) / student_count) if student_count else 0.0
    return {
        "confusion_count": confusion["level_percent"],
        "confusion_level_percent": confusion["level_percent"],
        "confusion_active_students": confusion["active_students"],
        "break_votes": len(session.break_votes),
        "student_count": student_count,
        "break_ratio": round(ratio, 3),
    }


def session_state_payload(session: RuntimeSession) -> dict[str, Any]:
    quiz_state_dict: dict[str, Any] = {
        "hidden": session.quiz_hidden,
        "cover_mode": session.quiz_cover_mode,
        "voting_closed": session.quiz_voting_closed,
        "answer_revealed": session.quiz_answer_revealed,
    }
    if session.quiz_answer_revealed and session.current_quiz:
        quiz_state_dict["correct_option_id"] = session.current_quiz.correct_option_id
        total = len(session.quiz_answers)
        per_option: dict[str, dict[str, Any]] = {}
        for opt in session.current_quiz.options:
            count = sum(1 for v in session.quiz_answers.values() if v == opt.id)
            per_option[opt.id] = {
                "count": count,
                "pct": round(count / total, 3) if total else 0.0,
            }
        quiz_state_dict["per_option"] = per_option
    return {
        "notes": session.notes,
        "break_active_until": session.break_active_until,
        "focus_period_ends_at": session.focus_period_ends_at,
        "quiz": session.current_quiz.model_dump() if session.current_quiz else None,
        "quiz_state": quiz_state_dict,
        "metrics": metrics_payload(session),
    }


def anonymous_questions_payload(session: RuntimeSession) -> dict[str, Any]:
    questions = [
        {
            "id": question.id,
            "text": question.text,
            "created_at": question.created_at,
            "resolved": question.resolved,
            "resolved_at": question.resolved_at,
        }
        for question in session.anonymous_questions
    ]
    pending_count = sum(1 for question in session.anonymous_questions if not question.resolved)
    return {
        "questions": questions,
        "pending_count": pending_count,
    }


SESSIONS: dict[str, RuntimeSession] = {}
BREAK_COOLDOWN_SECONDS = 30
BREAK_THRESHOLD_PERCENT = 0.4
MAX_BREAK_DURATION_SECONDS = 3600


def generate_session_code() -> str:
    alphabet = string.ascii_uppercase + string.digits
    while True:
        code = "".join(random.choices(alphabet, k=6))
        if code not in SESSIONS and not session_exists(code):
            return code


def get_teacher(session: RuntimeSession) -> ClientState | None:
    for c in session.clients.values():
        if c.role == "teacher":
            return c
    return None


async def send_json(ws: WebSocket, payload: dict[str, Any]) -> None:
    await ws.send_text(json.dumps(payload))


async def broadcast(session: RuntimeSession, payload: dict[str, Any], *, role: str | None = None) -> None:
    dead_clients: list[str] = []
    for client_id, client in session.clients.items():
        if role is not None and client.role != role:
            continue
        try:
            await send_json(client.websocket, payload)
        except Exception:
            dead_clients.append(client_id)

    for dead in dead_clients:
        session.clients.pop(dead, None)


def build_quiz_fallback(notes: str) -> QuizPayload:
    topic = (notes or "the current lecture").strip()
    if len(topic) > 70:
        topic = topic[:70] + "..."

    correct_id = "B"
    return QuizPayload(
        question=f"Which statement best summarizes {topic}?",
        options=[
            QuizOption(id="A", text="It is unrelated to the lesson topic."),
            QuizOption(id="B", text="It captures a key concept from the current lecture."),
            QuizOption(id="C", text="It should be ignored because it has no practical use."),
            QuizOption(id="D", text="It only applies to historical systems and not today."),
        ],
        correct_option_id=correct_id,
    )


QUIZ_PRESET_INSTRUCTIONS = {
    "default": "Create a clear concept-check question focused on the most important idea from the current context.",
    "funny": "Write a light, classroom-safe question with a subtle playful tone, while staying educational and accurate.",
    "challenge": "Write a challenging question that requires deeper reasoning, not just recall.",
    "misconception": "Target a common misconception and use plausible distractors that reveal misunderstanding.",
    "real_world": "Frame the question around a practical real-world scenario or application.",
}


def compose_quiz_generation_prompt(style_preset: str, custom_prompt: str) -> str:
    normalized_preset = (style_preset or "default").strip().lower()
    style_instruction = QUIZ_PRESET_INSTRUCTIONS.get(
        normalized_preset,
        QUIZ_PRESET_INSTRUCTIONS["default"],
    )
    custom_instruction = custom_prompt.strip()

    prompt = (
        "Generate exactly one multiple choice question with 4 options (A, B, C, D) based on lecture notes. "
        "Use the selected style instruction below. "
        "Keep all options similarly sized and similarly specific, so the correct answer is NOT obvious from option length, detail level, wording style, or formatting cues. "
        "Distractors must be plausible and non-joke unless style explicitly allows a playful tone. "
        "Do not add explanations.")

    prompt += f"\n\nStyle instruction:\n{style_instruction}"
    if custom_instruction:
        prompt += f"\n\nAdditional teacher instruction:\n{custom_instruction}"

    prompt += (
        "\n\nReturn JSON only with keys: question, options, correct_option_id. "
        "options must be an array of exactly 4 objects with keys id and text, and ids must be A, B, C, D. "
        "correct_option_id must be one of A, B, C, D."
    )
    return prompt


def build_quiz_with_ai(
    notes: str,
    style_preset: str = "default",
    custom_prompt: str = "",
) -> QuizPayload:
    _cache_key = ai_cache_key("quiz", notes or "", style_preset, custom_prompt or "")
    _cached = ai_cache_get(_cache_key)
    if _cached is not None:
        return QuizPayload(**json.loads(_cached))

    gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
    gemini_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
    errors: list[str] = []

    if gemini_api_key:
        prompt = compose_quiz_generation_prompt(style_preset, custom_prompt)

        parts: list[dict[str, Any]] = [
            {
                "text": f"Lecture notes:\n{notes or 'No notes provided.'}\n\n{prompt}",
            }
        ]

        models_to_try = [gemini_model]
        if gemini_model != "gemini-2.5-flash":
            models_to_try.append("gemini-2.5-flash")

        for model_name in models_to_try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={gemini_api_key}"
            body = {
                "contents": [{"role": "user", "parts": parts}],
                "generationConfig": {
                    "temperature": 0.2,
                    "responseMimeType": "application/json",
                },
            }

            try:
                req = Request(
                    url,
                    data=json.dumps(body).encode("utf-8"),
                    headers={"Content-Type": "application/json"},
                    method="POST",
                )
                with urlopen(req, timeout=20) as resp:
                    response_text = resp.read().decode("utf-8")
                response_json = json.loads(response_text)
                text_output = (
                    response_json.get("candidates", [{}])[0]
                    .get("content", {})
                    .get("parts", [{}])[0]
                    .get("text", "")
                    .strip()
                )
                if not text_output:
                    raise ValueError("Gemini returned an empty response")

                if text_output.startswith("```"):
                    text_output = text_output.strip("`")
                    if text_output.lower().startswith("json"):
                        text_output = text_output[4:].strip()

                parsed = json.loads(text_output)
                quiz = QuizPayload(**parsed)
                ai_cache_set(_cache_key, json.dumps(parsed), ttl_seconds=3600)
                return quiz
            except HTTPError as exc:
                details = ""
                try:
                    details = exc.read().decode("utf-8", errors="ignore")
                except Exception:
                    details = ""
                details = (details or exc.reason or "request failed")
                errors.append(f"Gemini model {model_name} HTTP {exc.code}: {details[:300]}")
            except (URLError, TimeoutError, json.JSONDecodeError, KeyError, ValueError) as exc:
                errors.append(f"Gemini model {model_name} error: {str(exc)[:300]}")

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    base_url = os.getenv("OPENAI_BASE_URL", "").strip() or None

    if api_key and OpenAI is not None:
        client = OpenAI(api_key=api_key, base_url=base_url)

        prompt = compose_quiz_generation_prompt(style_preset, custom_prompt)

        user_content: list[dict[str, Any]] = [
            {
                "type": "input_text",
                "text": f"Lecture notes:\n{notes or 'No notes provided.'}\n\n{prompt}",
            }
        ]
        try:
            response = client.responses.create(
                model=model,
                input=[
                    {
                        "role": "system",
                        "content": "You produce concise, educational quizzes in strict JSON.",
                    },
                    {
                        "role": "user",
                        "content": user_content,
                    },
                ],
                temperature=0.2,
            )

            text_output = response.output_text.strip()
            if not text_output:
                raise ValueError("OpenAI-compatible provider returned an empty response")

            if text_output.startswith("```"):
                text_output = text_output.strip("`")
                if text_output.lower().startswith("json"):
                    text_output = text_output[4:].strip()

            parsed = json.loads(text_output)
            quiz = QuizPayload(**parsed)
            ai_cache_set(_cache_key, json.dumps(parsed), ttl_seconds=3600)
            return quiz
        except Exception as exc:
            errors.append(f"OpenAI-compatible error: {str(exc)[:300]}")
    elif api_key and OpenAI is None:
        errors.append("OpenAI-compatible API key is set but OpenAI SDK is unavailable")

    if not gemini_api_key and not api_key:
        raise RuntimeError("No AI provider configured. Set GEMINI_API_KEY (recommended).")

    if errors:
        raise RuntimeError("; ".join(errors))

    raise RuntimeError("Quiz generation failed with unknown AI error")


def build_screen_explanation_with_ai(notes: str) -> str:
    _cache_key = ai_cache_key("explain", notes or "")
    _cached = ai_cache_get(_cache_key)
    if _cached is not None:
        return _cached

    gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
    gemini_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
    errors: list[str] = []

    if gemini_api_key:
        prompt = (
            "Explain what the teacher is currently showing in simple student-friendly language. "
            "Keep it concise (about 4-7 sentences), include the likely goal of the slide/screen, "
            "and suggest one practical thing the student should focus on next."
        )

        parts: list[dict[str, Any]] = [
            {
                "text": f"Lecture notes:\n{notes or 'No notes provided.'}\n\n{prompt}",
            }
        ]

        models_to_try = [gemini_model]
        if gemini_model != "gemini-2.5-flash":
            models_to_try.append("gemini-2.5-flash")

        for model_name in models_to_try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={gemini_api_key}"
            body = {
                "contents": [{"role": "user", "parts": parts}],
                "generationConfig": {
                    "temperature": 0.2,
                },
            }

            try:
                req = Request(
                    url,
                    data=json.dumps(body).encode("utf-8"),
                    headers={"Content-Type": "application/json"},
                    method="POST",
                )
                with urlopen(req, timeout=20) as resp:
                    response_text = resp.read().decode("utf-8")
                response_json = json.loads(response_text)
                text_output = (
                    response_json.get("candidates", [{}])[0]
                    .get("content", {})
                    .get("parts", [{}])[0]
                    .get("text", "")
                    .strip()
                )
                if text_output:
                    result = text_output[:1400]
                    ai_cache_set(_cache_key, result, ttl_seconds=3600)
                    return result
                raise ValueError("Gemini returned an empty explanation")
            except HTTPError as exc:
                details = ""
                try:
                    details = exc.read().decode("utf-8", errors="ignore")
                except Exception:
                    details = ""
                details = (details or exc.reason or "request failed")
                errors.append(f"Gemini model {model_name} HTTP {exc.code}: {details[:300]}")
            except (URLError, TimeoutError, json.JSONDecodeError, KeyError, ValueError) as exc:
                errors.append(f"Gemini model {model_name} error: {str(exc)[:300]}")

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    base_url = os.getenv("OPENAI_BASE_URL", "").strip() or None

    if api_key and OpenAI is not None:
        client = OpenAI(api_key=api_key, base_url=base_url)

        user_content: list[dict[str, Any]] = [
            {
                "type": "input_text",
                "text": (
                    "Explain what the teacher is currently showing in simple student-friendly language. "
                    "Keep it concise (about 4-7 sentences), include the likely goal of the slide/screen, "
                    "and suggest one practical thing the student should focus on next.\n\n"
                    f"Lecture notes:\n{notes or 'No notes provided.'}"
                ),
            }
        ]
        try:
            response = client.responses.create(
                model=model,
                input=[
                    {
                        "role": "system",
                        "content": "You explain live lecture visuals in concise and supportive language.",
                    },
                    {
                        "role": "user",
                        "content": user_content,
                    },
                ],
                temperature=0.2,
            )

            text_output = response.output_text.strip()
            if text_output:
                result = text_output[:1400]
                ai_cache_set(_cache_key, result, ttl_seconds=3600)
                return result
            raise ValueError("OpenAI-compatible provider returned an empty explanation")
        except Exception as exc:
            errors.append(f"OpenAI-compatible error: {str(exc)[:300]}")
    elif api_key and OpenAI is None:
        errors.append("OpenAI-compatible API key is set but OpenAI SDK is unavailable")

    if not gemini_api_key and not api_key:
        raise RuntimeError("No AI provider configured. Set GEMINI_API_KEY (recommended).")

    if errors:
        raise RuntimeError("; ".join(errors))

    raise RuntimeError("Screen explanation failed with unknown AI error")


def extract_text_from_presentation(file_path: Path, original_name: str, mime_type: str | None) -> str:
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md", ".csv", ".json"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")[:50000]

    if suffix == ".pdf" or (mime_type or "").lower() == "application/pdf":
        try:
            from pypdf import PdfReader
        except Exception as exc:
            raise RuntimeError("PDF parsing is unavailable. Install pypdf.") from exc

        reader = PdfReader(str(file_path))
        chunks: list[str] = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text)
            if sum(len(item) for item in chunks) > 50000:
                break
        return "\n\n".join(chunks)[:50000]

    if suffix == ".pptx" or "presentation" in (mime_type or "").lower():
        try:
            from pptx import Presentation
        except Exception as exc:
            raise RuntimeError("PPTX parsing is unavailable. Install python-pptx.") from exc

        prs = Presentation(str(file_path))
        chunks: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_lines.append(text.strip())
            if slide_lines:
                chunks.append(f"Slide {idx}:\n" + "\n".join(slide_lines))
            if sum(len(item) for item in chunks) > 50000:
                break
        return "\n\n".join(chunks)[:50000]

    raise RuntimeError(
        f"Unsupported presentation format for AI notes: {original_name}. Supported: .pptx, .pdf, .txt, .md, .csv, .json"
    )


def build_student_notes_with_ai(presentation_text: str, presentation_name: str) -> str:
    source = (presentation_text or "").strip()
    if not source:
        raise RuntimeError("No readable text was found in the presentation.")

    _cache_key = ai_cache_key("student_notes", source[:24000])
    _cached = ai_cache_get(_cache_key)
    if _cached is not None:
        return _cached

    source_excerpt = source[:24000]
    topic = Path(presentation_name).stem.replace("-", " ").replace("_", " ").strip() or "Science Topic"
    prompt = (
        "A hyper-realistic, top-down photograph of a single page of handwritten notes on clean white dotted journal paper, "
        "representing a \"modern minimalist\" aesthetic. "
        f"The subject is {topic}. "
        "The title is neatly lettered in fine black ink. "
        "The content is organized with thin black fineliner pen, featuring extremely precise, small print handwriting, "
        "simple bullet points, and neat numbered lists. "
        "Key terms are subtly highlighted with a single, very light pastel color (e.g., pale mint green). "
        "The notes are sparse, utilizing significant negative space. "
        "There is one simple, flawlessly executed black ink diagram with elegant, thin arrows. "
        "A single black fineliner pen rests diagonally beside the paper. "
        "The background is a plain, light wood grain desk surface. Natural, soft daylight. "
        "8k, macro detail, perfect legibility, high resolution. "
        "\n\nNow generate ONLY the note content text that should be written on that page. "
        "Do not describe camera, paper, desk, lighting, or photo style. "
        "Output plain text only (no markdown, no code fences). "
        "Keep it concise and student-friendly so it fits one page."
    )

    gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
    gemini_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")
    errors: list[str] = []

    if gemini_api_key:
        parts: list[dict[str, Any]] = [
            {
                "text": f"Presentation file: {presentation_name}\n\nContent:\n{source_excerpt}\n\n{prompt}",
            }
        ]

        models_to_try = [gemini_model]
        if gemini_model != "gemini-2.5-flash":
            models_to_try.append("gemini-2.5-flash")

        for model_name in models_to_try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={gemini_api_key}"
            body = {
                "contents": [{"role": "user", "parts": parts}],
                "generationConfig": {
                    "temperature": 0.2,
                },
            }
            try:
                req = Request(
                    url,
                    data=json.dumps(body).encode("utf-8"),
                    headers={"Content-Type": "application/json"},
                    method="POST",
                )
                with urlopen(req, timeout=30) as resp:
                    response_text = resp.read().decode("utf-8")
                response_json = json.loads(response_text)
                text_output = (
                    response_json.get("candidates", [{}])[0]
                    .get("content", {})
                    .get("parts", [{}])[0]
                    .get("text", "")
                    .strip()
                )
                if text_output:
                    result = text_output[:6000]
                    ai_cache_set(_cache_key, result, ttl_seconds=604800)
                    return result
                raise ValueError("Gemini returned empty notes")
            except HTTPError as exc:
                details = ""
                try:
                    details = exc.read().decode("utf-8", errors="ignore")
                except Exception:
                    details = ""
                details = (details or exc.reason or "request failed")
                errors.append(f"Gemini model {model_name} HTTP {exc.code}: {details[:300]}")
            except (URLError, TimeoutError, json.JSONDecodeError, KeyError, ValueError) as exc:
                errors.append(f"Gemini model {model_name} error: {str(exc)[:300]}")

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    base_url = os.getenv("OPENAI_BASE_URL", "").strip() or None

    if api_key and OpenAI is not None:
        client = OpenAI(api_key=api_key, base_url=base_url)
        try:
            response = client.responses.create(
                model=model,
                input=[
                    {
                        "role": "system",
                        "content": "You write clear, student-friendly study notes.",
                    },
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "input_text",
                                "text": f"Presentation file: {presentation_name}\n\nContent:\n{source_excerpt}\n\n{prompt}",
                            }
                        ],
                    },
                ],
                temperature=0.2,
            )
            text_output = response.output_text.strip()
            if text_output:
                result = text_output[:6000]
                ai_cache_set(_cache_key, result, ttl_seconds=604800)
                return result
            raise ValueError("OpenAI-compatible provider returned empty notes")
        except Exception as exc:
            errors.append(f"OpenAI-compatible error: {str(exc)[:300]}")
    elif api_key and OpenAI is None:
        errors.append("OpenAI-compatible API key is set but OpenAI SDK is unavailable")

    if not gemini_api_key and not api_key:
        raise RuntimeError("No AI provider configured. Set GEMINI_API_KEY (recommended).")

    if errors:
        raise RuntimeError("; ".join(errors))

    raise RuntimeError("Student-notes generation failed with unknown AI error")


STUDENT_NOTES_IMAGE_PROMPT_TEMPLATE = (
    "A high-resolution, top-down flat lay of a single page of handwritten student notes on white "
    "[lined/grid] paper, isolated on a plain background for easy PNG conversion. The style is \"Gen Z studygram\": "
    "neat but organic handwriting using a black 0.5mm gel pen, with key terms highlighted in mild pastel colors "
    "(lemon yellow and mint green). Includes hand-drawn arrows, small aesthetic doodles in the margins "
    "(like tiny stars or chemistry beakers), and a \"header\" title in simple faux-calligraphy. "
    "The paper has slight natural crinkles. Lighting is bright, even, and clinical. Soft shadows only. "
    "8k resolution, macro photography style, extremely legible text."
)


def build_notes_image_prompt(notes_text: str) -> str:
    normalized_notes = (notes_text or "").strip()
    return (
        f"{STUDENT_NOTES_IMAGE_PROMPT_TEMPLATE}\n\n"
        "Write exactly this study-notes content on the page and keep it fully legible:\n"
        "---\n"
        f"{normalized_notes[:5000]}\n"
        "---\n"
        "Do not add logos, watermarks, or extra pages. Keep one page only."
    )


def generate_notes_png_with_ai(notes_text: str) -> bytes:
    gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
    configured_model = os.getenv("GEMINI_IMAGE_MODEL", "").strip()

    if not gemini_api_key:
        raise RuntimeError("Gemini image generation is unavailable. Set GEMINI_API_KEY.")

    prompt = build_notes_image_prompt(notes_text)

    models_to_try: list[str] = []

    if configured_model:
        models_to_try.append(configured_model)

    # Try likely image-capable model aliases first.
    for candidate in [
        "gemini-2.5-flash-image",
        "gemini-2.5-flash-image-preview",
        "gemini-2.0-flash-preview-image-generation",
        "gemini-2.0-flash-exp-image-generation",
    ]:
        if candidate not in models_to_try:
            models_to_try.append(candidate)

    # Discover account-available models and prioritize image-related ones.
    try:
        list_url = f"https://generativelanguage.googleapis.com/v1beta/models?key={gemini_api_key}"
        list_req = Request(list_url, headers={"Content-Type": "application/json"}, method="GET")
        with urlopen(list_req, timeout=20) as list_resp:
            list_text = list_resp.read().decode("utf-8")
        list_json = json.loads(list_text)
        discovered = list_json.get("models", [])
        for model_entry in discovered:
            name = str(model_entry.get("name", ""))
            if not name:
                continue
            short_name = name.split("models/", 1)[-1]
            generation_methods = [str(item) for item in model_entry.get("supportedGenerationMethods", [])]
            if "generateContent" not in generation_methods:
                continue

            lower = short_name.lower()
            # Heuristic: image/vision/preview variants are more likely to return inline images.
            if any(token in lower for token in ["image", "vision", "imagen", "preview"]):
                if short_name not in models_to_try:
                    models_to_try.append(short_name)
    except Exception:
        # Model discovery is best-effort; keep static fallback list.
        pass

    errors: list[str] = []

    def decode_image_from_response(response_json: dict[str, Any]) -> bytes | None:
        candidates = response_json.get("candidates", [])
        for candidate in candidates:
            parts = candidate.get("content", {}).get("parts", [])
            for part in parts:
                inline = part.get("inlineData") or part.get("inline_data")
                if not inline:
                    continue
                mime_type = (inline.get("mimeType") or inline.get("mime_type") or "").lower()
                if mime_type not in {"image/png", "image/jpeg", "image/webp"}:
                    continue
                data = inline.get("data")
                if data:
                    return b64decode(data)
        return None

    for image_model in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{image_model}:generateContent?key={gemini_api_key}"
            body = {
                "contents": [
                    {
                        "role": "user",
                        "parts": [{"text": prompt}],
                    }
                ],
                "generationConfig": {
                    "temperature": 0.35,
                    "responseModalities": ["IMAGE", "TEXT"],
                },
            }

            req = Request(
                url,
                data=json.dumps(body).encode("utf-8"),
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urlopen(req, timeout=80) as resp:
                response_text = resp.read().decode("utf-8")

            response_json = json.loads(response_text)
            image_bytes = decode_image_from_response(response_json)
            if image_bytes:
                return image_bytes
            errors.append(f"{image_model}: response had no inline image bytes")
        except HTTPError as exc:
            details = ""
            try:
                details = exc.read().decode("utf-8", errors="ignore")
            except Exception:
                details = ""
            details = details or exc.reason or "request failed"
            errors.append(f"{image_model} HTTP {exc.code}: {details[:240]}")
        except Exception as exc:
            errors.append(f"{image_model}: {str(exc)[:240]}")

    raise RuntimeError("Gemini image generation failed: " + " | ".join(errors[:3]))


def render_notes_png(title: str, notes_text: str) -> bytes:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception as exc:
        raise RuntimeError("PNG rendering is unavailable. Install Pillow.") from exc

    width = 1600
    height = 2200
    margin = 90

    image = Image.new("RGB", (width, height), color=(251, 252, 255))
    draw = ImageDraw.Draw(image)

    try:
        title_font = ImageFont.truetype("arial.ttf", 54)
        body_font = ImageFont.truetype("arial.ttf", 30)
    except Exception:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    current_y = margin
    draw.text((margin, current_y), title, fill=(20, 34, 62), font=title_font)
    current_y += 90
    draw.line((margin, current_y, width - margin, current_y), fill=(186, 200, 230), width=3)
    current_y += 40

    usable_width = width - (2 * margin)
    paragraphs = [line.strip() for line in notes_text.splitlines() if line.strip()]
    if not paragraphs:
        paragraphs = ["No notes could be generated for this presentation."]

    rendered_lines: list[str] = []
    for paragraph in paragraphs:
        wrapped = textwrap.wrap(paragraph, width=78) or [paragraph]
        rendered_lines.extend(wrapped)
        rendered_lines.append("")

    max_lines = 58
    if len(rendered_lines) > max_lines:
        rendered_lines = rendered_lines[: max_lines - 1] + ["... (content truncated)"]

    for line in rendered_lines:
        if line == "":
            current_y += 16
            continue

        while draw.textlength(line, font=body_font) > usable_width and len(line) > 4:
            line = line[:-2]
        draw.text((margin, current_y), line, fill=(38, 46, 66), font=body_font)
        current_y += 42
        if current_y > height - margin:
            break

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def get_accessible_presentation_row(
    cursor: sqlite3.Cursor,
    *,
    user: dict[str, Any],
    presentation_id: int,
    normalized_session_code: str | None,
) -> sqlite3.Row | None:
    if normalized_session_code:
        cursor.execute("SELECT teacher_name, owner_user_id FROM sessions WHERE code = ?", (normalized_session_code,))
        session_row = cursor.fetchone()
        if not session_row:
            raise HTTPException(status_code=404, detail="Session not found")

        session_owner_id = session_row["owner_user_id"]
        is_session_owner = (
            (session_owner_id is not None and session_owner_id == user["id"])
            or (session_owner_id is None and session_row["teacher_name"] == user["display_name"])
        )
        if is_session_owner:
            cursor.execute(
                """
                SELECT id, session_code, original_name, stored_name, mime_type
                FROM presentations
                WHERE id = ? AND user_id = ?
                """,
                (presentation_id, user["id"]),
            )
        else:
            cursor.execute(
                """
                SELECT p.id, p.session_code, p.original_name, p.stored_name, p.mime_type
                FROM presentations p
                JOIN sessions s ON s.code = p.session_code
                JOIN users u ON u.id = p.user_id
                WHERE p.id = ?
                  AND p.session_code = ?
                                    AND (
                                        (s.owner_user_id IS NOT NULL AND s.owner_user_id = u.id)
                                        OR (s.owner_user_id IS NULL AND s.teacher_name = u.display_name)
                                    )
                """,
                (presentation_id, normalized_session_code),
            )
    else:
        cursor.execute(
            """
            SELECT id, session_code, original_name, stored_name, mime_type
            FROM presentations
            WHERE id = ? AND user_id = ?
            """,
            (presentation_id, user["id"]),
        )
    return cursor.fetchone()


def analytics_for_session(session: RuntimeSession) -> dict[str, Any]:
    confusion = confusion_snapshot(session)
    answers = session.quiz_answers.values()
    total_answers = len(session.quiz_answers)
    correct_answers = 0
    if session.current_quiz:
        correct_answers = sum(1 for answer in answers if answer == session.current_quiz.correct_option_id)

    accuracy = (correct_answers / total_answers) if total_answers else 0.0
    active_students = confusion["student_count"]
    break_vote_rate = (len(session.break_votes) / active_students) if active_students else 0.0
    quiz_participation_rate = (total_answers / active_students) if active_students else 0.0
    confusion_per_student = confusion["average_level"] if active_students else 0.0

    score_raw = (
        min(quiz_participation_rate, 1.0) * 0.55
        + min(break_vote_rate / BREAK_THRESHOLD_PERCENT, 1.0) * 0.25
        + min(confusion_per_student, 1.0) * 0.20
    )
    engagement_score = int(round(score_raw * 100))

    end_epoch = session.ended_at_epoch if session.ended_at_epoch is not None else time.time()
    duration_seconds = max(0, int(end_epoch - session.created_at_epoch))

    return {
        "session_code": session.code,
        "teacher_name": session.teacher_name,
        "session_active": session.active,
        "duration_seconds": duration_seconds,
        "student_count": active_students,
        "confusion_count": confusion["level_percent"],
        "confusion_level_percent": confusion["level_percent"],
        "confusion_active_students": confusion["active_students"],
        "break_votes": len(session.break_votes),
        "break_active": bool(session.break_active_until and session.break_active_until > time.time()),
        "quiz": {
            "total_answers": total_answers,
            "correct_answers": correct_answers,
            "accuracy": round(accuracy, 3),
            "hidden": session.quiz_hidden,
            "cover_mode": session.quiz_cover_mode,
            "voting_closed": session.quiz_voting_closed,
            "answer_revealed": session.quiz_answer_revealed,
        },
        "engagement": {
            "score": engagement_score,
            "quiz_participation_rate": round(quiz_participation_rate, 3),
            "break_vote_rate": round(break_vote_rate, 3),
            "confusion_per_student": round(confusion_per_student, 3),
        },
        "notes": session.notes,
    }


def record_engagement_point(session: RuntimeSession, source: str) -> None:
    analytics = analytics_for_session(session)
    point = {
        "recorded_at_epoch": round(time.time(), 3),
        "source": source,
        "engagement_score": int(analytics.get("engagement", {}).get("score", 0)),
        "confusion_level_percent": int(analytics.get("confusion_level_percent", 0)),
        "break_votes": int(analytics.get("break_votes", 0)),
    }

    timeline = session.engagement_timeline
    if timeline:
        last = timeline[-1]
        if (
            last.get("engagement_score") == point["engagement_score"]
            and last.get("confusion_level_percent") == point["confusion_level_percent"]
            and last.get("break_votes") == point["break_votes"]
            and last.get("source") == point["source"]
            and (point["recorded_at_epoch"] - float(last.get("recorded_at_epoch", 0))) < 5
        ):
            return

    timeline.append(point)


def build_full_analytics_report(session: RuntimeSession) -> dict[str, Any]:
    analytics = analytics_for_session(session)
    students_connected = []
    now_epoch = time.time()
    for client in session.clients.values():
        if client.role != "student":
            continue
        students_connected.append(
            {
                "client_id": client.client_id,
                "name": client.name,
                "joined_at_epoch": round(client.joined_at, 3),
                "time_in_session_seconds": max(0, int(now_epoch - client.joined_at)),
            }
        )

    timeline = sorted(
        session.engagement_timeline,
        key=lambda item: float(item.get("recorded_at_epoch", 0)),
    )
    if not timeline:
        record_engagement_point(session, "report_generated")
        timeline = sorted(
            session.engagement_timeline,
            key=lambda item: float(item.get("recorded_at_epoch", 0)),
        )

    return {
        "report_type": "session_engagement",
        "generated_at": now_iso(),
        "analytics": analytics,
        "students_connected": students_connected,
        "engagement_timeline": timeline,
    }


def _normalize_insights_payload(payload: dict[str, Any]) -> dict[str, Any]:
    def clean_text(value: Any, *, fallback: str, max_len: int = 700) -> str:
        text = " ".join(str(value or "").split())
        if not text:
            text = fallback
        return text[:max_len]

    def clean_list(values: Any, *, fallback: list[str], max_items: int = 5) -> list[str]:
        if not isinstance(values, list):
            values = []
        cleaned: list[str] = []
        for item in values:
            text = " ".join(str(item or "").split())
            if text:
                cleaned.append(text[:220])
            if len(cleaned) >= max_items:
                break
        return cleaned or fallback

    return {
        "title": clean_text(payload.get("title"), fallback="Session Analytics Insights", max_len=120),
        "executive_summary": clean_text(
            payload.get("executive_summary"),
            fallback="Class engagement summary generated from session analytics.",
        ),
        "key_findings": clean_list(
            payload.get("key_findings"),
            fallback=["Engagement indicators were collected successfully for this session."],
        ),
        "risks": clean_list(
            payload.get("risks"),
            fallback=["No critical risk was detected from the available analytics."],
        ),
        "recommendations": clean_list(
            payload.get("recommendations"),
            fallback=["Continue tracking participation and confusion trends across sessions."],
        ),
    }


def build_analytics_insights_fallback(report: dict[str, Any]) -> dict[str, Any]:
    analytics = report.get("analytics", {})
    engagement = analytics.get("engagement", {})
    quiz = analytics.get("quiz", {})

    score = int(engagement.get("score") or 0)
    student_count = int(analytics.get("student_count") or 0)
    confusion_percent = int(analytics.get("confusion_level_percent") or 0)
    break_votes = int(analytics.get("break_votes") or 0)
    participation_rate = float(engagement.get("quiz_participation_rate") or 0.0)
    accuracy = float(quiz.get("accuracy") or 0.0)

    if score >= 75:
        engagement_label = "high"
    elif score >= 45:
        engagement_label = "moderate"
    else:
        engagement_label = "low"

    risks: list[str] = []
    if participation_rate < 0.55:
        risks.append("Quiz participation was below 55%, indicating uneven student involvement.")
    if accuracy < 0.5:
        risks.append("Quiz accuracy was below 50%, suggesting students may need concept reinforcement.")
    if confusion_percent >= 60:
        risks.append("High confusion levels were detected, which can reduce learning retention.")
    if break_votes >= max(2, int(student_count * 0.35)):
        risks.append("Break demand was elevated, which may indicate cognitive overload or pacing issues.")
    if not risks:
        risks.append("No major risk trend was identified from the available analytics.")

    recommendations = [
        "Start the next class with a short recap of the most difficult concept from this session.",
        "Add one low-stakes check-in question every 8-12 minutes to improve participation consistency.",
        "Trigger a micro-break earlier when break-vote and confusion indicators begin rising together.",
    ]

    return _normalize_insights_payload(
        {
            "title": "AI-Assisted Session Insights",
            "executive_summary": (
                f"This session shows {engagement_label} engagement (score {score}/100) across {student_count} active students. "
                f"Quiz participation was {round(participation_rate * 100)}% with {round(accuracy * 100)}% accuracy, "
                f"while confusion reached {confusion_percent}% and break votes totaled {break_votes}."
            ),
            "key_findings": [
                f"Engagement score reached {score}/100.",
                f"Quiz participation rate was {round(participation_rate * 100)}%.",
                f"Quiz accuracy ended at {round(accuracy * 100)}%.",
                f"Confusion level indicator was {confusion_percent}%.",
            ],
            "risks": risks,
            "recommendations": recommendations,
        }
    )


def build_analytics_insights_with_ai(report: dict[str, Any]) -> dict[str, Any]:
    analytics = report.get("analytics", {})
    prompt = (
        "You are an educational analytics assistant. "
        "Generate concise insights from classroom session metrics. "
        "Return STRICT JSON only with keys: title, executive_summary, key_findings, risks, recommendations. "
        "key_findings, risks, recommendations must each be arrays of short bullet-style strings. "
        "Focus on actionable and evidence-based insights."
    )
    input_data = {
        "session_code": analytics.get("session_code"),
        "teacher_name": analytics.get("teacher_name"),
        "duration_seconds": analytics.get("duration_seconds"),
        "student_count": analytics.get("student_count"),
        "confusion_level_percent": analytics.get("confusion_level_percent"),
        "break_votes": analytics.get("break_votes"),
        "quiz": analytics.get("quiz"),
        "engagement": analytics.get("engagement"),
    }

    gemini_api_key = os.getenv("GEMINI_API_KEY", "").strip()
    gemini_model = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")

    if gemini_api_key:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{gemini_model}:generateContent?key={gemini_api_key}"
            body = {
                "contents": [
                    {
                        "role": "user",
                        "parts": [
                            {
                                "text": f"Session analytics data:\n{json.dumps(input_data, ensure_ascii=True)}\n\n{prompt}",
                            }
                        ],
                    }
                ],
                "generationConfig": {
                    "temperature": 0.2,
                    "responseMimeType": "application/json",
                },
            }
            req = Request(
                url,
                data=json.dumps(body).encode("utf-8"),
                headers={"Content-Type": "application/json"},
                method="POST",
            )
            with urlopen(req, timeout=20) as resp:
                response_text = resp.read().decode("utf-8")
            parsed_response = json.loads(response_text)
            model_text = (
                parsed_response.get("candidates", [{}])[0]
                .get("content", {})
                .get("parts", [{}])[0]
                .get("text", "")
                .strip()
            )
            if model_text.startswith("```"):
                model_text = model_text.strip("`")
                if model_text.lower().startswith("json"):
                    model_text = model_text[4:].strip()
            if model_text:
                return _normalize_insights_payload(json.loads(model_text))
        except Exception:
            pass

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    model = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
    base_url = os.getenv("OPENAI_BASE_URL", "").strip() or None
    if api_key and OpenAI is not None:
        try:
            client = OpenAI(api_key=api_key, base_url=base_url)
            response = client.responses.create(
                model=model,
                input=[
                    {
                        "role": "system",
                        "content": "You produce concise educational analytics insights in strict JSON.",
                    },
                    {
                        "role": "user",
                        "content": (
                            f"Session analytics data:\n{json.dumps(input_data, ensure_ascii=True)}\n\n{prompt}"
                        ),
                    },
                ],
                temperature=0.2,
            )
            text_output = (response.output_text or "").strip()
            if text_output.startswith("```"):
                text_output = text_output.strip("`")
                if text_output.lower().startswith("json"):
                    text_output = text_output[4:].strip()
            if text_output:
                return _normalize_insights_payload(json.loads(text_output))
        except Exception:
            pass

    return build_analytics_insights_fallback(report)


def _wrap_text_lines(text: str, font_name: str, font_size: int, max_width: float) -> list[str]:
    if not text:
        return [""]

    if stringWidth is None:
        return [text]

    words = text.split()
    if not words:
        return [""]

    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if stringWidth(candidate, font_name, font_size) <= max_width:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def generate_session_report_pdf(report: dict[str, Any], insights: dict[str, Any]) -> bytes:
    if canvas is None or A4 is None or cm is None:
        raise RuntimeError("PDF generation dependency is not available")

    analytics = report.get("analytics", {})
    engagement = analytics.get("engagement", {})
    quiz = analytics.get("quiz", {})
    students = report.get("students_connected", [])
    timeline = report.get("engagement_timeline", [])

    duration_seconds = int(analytics.get("duration_seconds") or 0)
    if not timeline:
        timeline = [
            {
                "recorded_at_epoch": 0,
                "engagement_score": int(engagement.get("score", 0)),
                "confusion_level_percent": int(analytics.get("confusion_level_percent", 0)),
            },
            {
                "recorded_at_epoch": max(duration_seconds, 1),
                "engagement_score": int(engagement.get("score", 0)),
                "confusion_level_percent": int(analytics.get("confusion_level_percent", 0)),
            },
        ]

    first_epoch = float(timeline[0].get("recorded_at_epoch", 0)) if timeline else 0.0
    line_points: list[dict[str, float]] = []
    for item in timeline:
        raw_epoch = float(item.get("recorded_at_epoch", first_epoch))
        elapsed = raw_epoch - first_epoch
        if duration_seconds > 0:
            elapsed = min(max(elapsed, 0.0), float(duration_seconds))
        line_points.append(
            {
                "elapsed": elapsed,
                "engagement": min(max(float(item.get("engagement_score", 0)), 0.0), 100.0),
                "confusion": min(max(float(item.get("confusion_level_percent", 0)), 0.0), 100.0),
            }
        )

    if len(line_points) == 1:
        line_points.append(
            {
                "elapsed": max(float(duration_seconds), 1.0),
                "engagement": line_points[0]["engagement"],
                "confusion": line_points[0]["confusion"],
            }
        )

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    page_width, page_height = A4

    margin_x = 2 * cm
    y = page_height - 2 * cm
    content_width = page_width - (2 * margin_x)
    bottom_limit = 2.1 * cm

    c_text = (0.15, 0.18, 0.24)
    c_muted = (0.46, 0.51, 0.60)
    c_header_bg = (0.92, 0.96, 1.0)
    c_card_bg = (0.97, 0.98, 1.0)
    c_card_border = (0.85, 0.89, 0.95)
    c_grid = (0.90, 0.93, 0.97)
    c_engagement = (0.07, 0.42, 0.73)
    c_confusion = (0.88, 0.45, 0.20)
    c_quiz_correct = (0.16, 0.62, 0.37)
    c_quiz_incorrect = (0.84, 0.33, 0.31)
    c_quiz_unanswered = (0.64, 0.66, 0.70)

    def new_page() -> None:
        nonlocal y
        pdf.showPage()
        y = page_height - 2 * cm

    def ensure_space(height_needed: float) -> None:
        nonlocal y
        if y - height_needed < bottom_limit:
            new_page()

    def draw_heading(
        text: str,
        *,
        size: int = 14,
        gap_before: int = 10,
        gap_after: int = 12,
    ) -> None:
        nonlocal y
        y -= gap_before
        ensure_space(size + gap_after + 6)
        pdf.setFont("Helvetica-Bold", size)
        pdf.setFillColorRGB(*c_text)
        pdf.drawString(margin_x, y, text)
        y -= gap_after + size

    def draw_paragraph(text: str, *, font: str = "Helvetica", size: int = 11, leading: int = 14) -> None:
        nonlocal y
        for paragraph in (text or "").splitlines() or [""]:
            lines = _wrap_text_lines(paragraph.strip(), font, size, content_width)
            for line in lines:
                ensure_space(leading + 2)
                pdf.setFont(font, size)
                pdf.setFillColorRGB(*c_text)
                pdf.drawString(margin_x, y, line)
                y -= leading
            y -= 2

    def draw_bullets(items: list[str]) -> None:
        nonlocal y
        for item in items:
            bullet_lines = _wrap_text_lines(item, "Helvetica", 11, content_width - 16)
            ensure_space((len(bullet_lines) + 1) * 14 + 4)
            pdf.setFont("Helvetica", 11)
            pdf.setFillColorRGB(*c_text)
            pdf.drawString(margin_x, y, "-")
            for line in bullet_lines:
                ensure_space(14)
                pdf.drawString(margin_x + 12, y, line)
                y -= 14
            y -= 2

    def draw_info_card(height: float) -> tuple[float, float, float, float]:
        nonlocal y
        ensure_space(height)
        top = y
        bottom = y - height
        pdf.setFillColorRGB(*c_card_bg)
        pdf.setStrokeColorRGB(*c_card_border)
        pdf.setLineWidth(1)
        pdf.roundRect(margin_x, bottom, content_width, height, 8, stroke=1, fill=1)
        y = bottom - 12
        return (margin_x, bottom, content_width, top)

    # Header band.
    header_height = 2.9 * cm
    ensure_space(header_height + 12)
    header_bottom = y - header_height
    pdf.setFillColorRGB(*c_header_bg)
    pdf.setStrokeColorRGB(*c_card_border)
    pdf.roundRect(margin_x, header_bottom, content_width, header_height, 10, stroke=1, fill=1)
    pdf.setFillColorRGB(*c_text)
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(margin_x + 12, y - 24, insights.get("title", "Session Analytics Report"))
    pdf.setFont("Helvetica", 9)
    pdf.setFillColorRGB(*c_muted)
    pdf.drawString(
        margin_x + 12,
        y - 40,
        (
            f"Session {analytics.get('session_code', 'N/A')} | "
            f"Teacher {analytics.get('teacher_name', 'N/A')} | "
            f"Generated {report.get('generated_at', now_iso())}"
        ),
    )
    y = header_bottom - 14

    draw_heading("Executive Summary", gap_before=2, gap_after=12)
    draw_paragraph(insights.get("executive_summary", "No executive summary available."))

    draw_heading("Core Metrics", gap_before=10, gap_after=12)
    card_x, card_bottom, card_width, card_top = draw_info_card(104)
    card_pad = 12
    col_w = (card_width - (card_pad * 2)) / 2
    metric_rows = [
        ("Duration", f"{analytics.get('duration_seconds', 0)}s"),
        ("Students", str(analytics.get("student_count", 0))),
        ("Engagement", f"{engagement.get('score', 0)} / 100"),
        ("Confusion", f"{analytics.get('confusion_level_percent', 0)}%"),
        ("Break votes", str(analytics.get("break_votes", 0))),
        (
            "Quiz accuracy",
            f"{round(float(quiz.get('accuracy', 0.0)) * 100)}% ({quiz.get('correct_answers', 0)}/{quiz.get('total_answers', 0)})",
        ),
    ]
    for idx, (label, value) in enumerate(metric_rows):
        col = idx % 2
        row = idx // 2
        x = card_x + card_pad + (col * col_w)
        row_y = card_top - 20 - (row * 28)
        pdf.setFont("Helvetica", 9)
        pdf.setFillColorRGB(*c_muted)
        pdf.drawString(x, row_y, label)
        pdf.setFont("Helvetica-Bold", 11)
        pdf.setFillColorRGB(*c_text)
        pdf.drawString(x, row_y - 12, value)

    draw_heading("Student Engagement Trend", gap_before=10, gap_after=12)
    chart_block_h = 188
    block_x, block_bottom, block_w, block_top = draw_info_card(chart_block_h)

    legend_y = block_top - 16
    pdf.setFillColorRGB(*c_engagement)
    pdf.rect(block_x + 12, legend_y - 4, 10, 6, stroke=0, fill=1)
    pdf.setFillColorRGB(*c_text)
    pdf.setFont("Helvetica", 9)
    pdf.drawString(block_x + 26, legend_y - 3, "Engagement")
    pdf.setFillColorRGB(*c_confusion)
    pdf.rect(block_x + 96, legend_y - 4, 10, 6, stroke=0, fill=1)
    pdf.setFillColorRGB(*c_text)
    pdf.drawString(block_x + 110, legend_y - 3, "Confusion")

    chart_left = block_x + 18
    chart_right = block_x + block_w - 14
    chart_width = chart_right - chart_left
    chart_bottom = block_bottom + 26
    chart_height = chart_block_h - 52

    pdf.setStrokeColorRGB(*c_grid)
    pdf.setLineWidth(0.6)
    for value in [0, 25, 50, 75, 100]:
        y_tick = chart_bottom + (chart_height * (value / 100.0))
        pdf.line(chart_left, y_tick, chart_right, y_tick)
        pdf.setFillColorRGB(*c_muted)
        pdf.setFont("Helvetica", 8)
        pdf.drawString(chart_left - 16, y_tick - 2, str(value))

    pdf.setStrokeColorRGB(0.72, 0.77, 0.84)
    pdf.setLineWidth(1)
    pdf.line(chart_left, chart_bottom, chart_left, chart_bottom + chart_height)
    pdf.line(chart_left, chart_bottom, chart_right, chart_bottom)

    x_max = max(float(duration_seconds), max((point["elapsed"] for point in line_points), default=1.0), 1.0)

    def draw_series(key: str, color: tuple[float, float, float]) -> None:
        pdf.setStrokeColorRGB(*color)
        pdf.setFillColorRGB(*color)
        pdf.setLineWidth(1.8)
        for idx in range(1, len(line_points)):
            p1 = line_points[idx - 1]
            p2 = line_points[idx]
            x1 = chart_left + (p1["elapsed"] / x_max) * chart_width
            y1 = chart_bottom + (p1[key] / 100.0) * chart_height
            x2 = chart_left + (p2["elapsed"] / x_max) * chart_width
            y2 = chart_bottom + (p2[key] / 100.0) * chart_height
            pdf.line(x1, y1, x2, y2)
        for point in line_points:
            x = chart_left + (point["elapsed"] / x_max) * chart_width
            y_point = chart_bottom + (point[key] / 100.0) * chart_height
            pdf.circle(x, y_point, 1.5, stroke=0, fill=1)

    draw_series("engagement", c_engagement)
    draw_series("confusion", c_confusion)

    pdf.setFillColorRGB(*c_muted)
    pdf.setFont("Helvetica", 8)
    pdf.drawString(chart_left, chart_bottom - 14, "Session timeline (seconds)")
    pdf.drawRightString(chart_right, chart_bottom - 14, f"{int(round(x_max))}s")

    draw_heading("Quiz Performance Visualization", gap_before=10, gap_after=12)
    quiz_block_h = 120
    q_x, q_bottom, q_w, q_top = draw_info_card(quiz_block_h)
    q_chart_left = q_x + 14
    q_chart_right = q_x + q_w - 14
    q_chart_w = q_chart_right - q_chart_left

    total_answers = int(quiz.get("total_answers", 0))
    correct_answers = int(quiz.get("correct_answers", 0))
    incorrect_answers = max(total_answers - correct_answers, 0)
    student_count = int(analytics.get("student_count", 0))
    unanswered = max(student_count - total_answers, 0)
    max_bar = max(student_count, total_answers, 1)

    quiz_bars = [
        ("Correct", correct_answers, c_quiz_correct),
        ("Incorrect", incorrect_answers, c_quiz_incorrect),
        ("Unanswered", unanswered, c_quiz_unanswered),
    ]

    for idx, (label, value, color) in enumerate(quiz_bars):
        row_y = q_top - 24 - (idx * 30)
        bar_w = (value / max_bar) * (q_chart_w - 88)
        pdf.setFillColorRGB(*c_muted)
        pdf.setFont("Helvetica", 9)
        pdf.drawString(q_chart_left, row_y, label)
        pdf.setFillColorRGB(0.92, 0.94, 0.97)
        pdf.roundRect(q_chart_left + 56, row_y - 7, q_chart_w - 88, 10, 3, stroke=0, fill=1)
        pdf.setFillColorRGB(*color)
        pdf.roundRect(q_chart_left + 56, row_y - 7, max(bar_w, 0.5), 10, 3, stroke=0, fill=1)
        pdf.setFillColorRGB(*c_text)
        pdf.setFont("Helvetica-Bold", 9)
        pdf.drawRightString(q_chart_right, row_y, str(value))

    accuracy_pct = round(float(quiz.get("accuracy", 0.0)) * 100)
    participation_pct = round((total_answers / max(student_count, 1)) * 100) if student_count else 0
    pdf.setFont("Helvetica", 9)
    pdf.setFillColorRGB(*c_muted)
    pdf.drawString(q_chart_left, q_bottom + 10, f"Accuracy: {accuracy_pct}%")
    pdf.drawString(q_chart_left + 112, q_bottom + 10, f"Participation: {participation_pct}%")

    draw_heading("Key Findings", gap_before=12, gap_after=12)
    draw_bullets(insights.get("key_findings", []))

    draw_heading("Risks", gap_before=12, gap_after=12)
    draw_bullets(insights.get("risks", []))

    draw_heading("Recommendations", gap_before=12, gap_after=12)
    draw_bullets(insights.get("recommendations", []))

    draw_heading("Connected Students", gap_before=12, gap_after=12)
    if students:
        for student in students[:25]:
            draw_paragraph(
                f"{student.get('name', 'Unknown')} (id: {student.get('client_id', 'n/a')}) - "
                f"time in session: {student.get('time_in_session_seconds', 0)}s",
                size=10,
                leading=13,
            )
    else:
        draw_paragraph("No student connection records were captured for this report.")

    pdf.save()
    return buffer.getvalue()


def current_assessment_question_id(session: RuntimeSession) -> str:
    if session.current_quiz_saved_id is not None:
        return f"quiz-{session.current_quiz_saved_id}"
    if session.current_quiz is None:
        return "quiz-unknown"
    digest = sha256(session.current_quiz.question.encode("utf-8")).hexdigest()[:12]
    return f"quiz-{digest}"


def ensure_assessment_schema() -> None:
    init_assessment_schema(DB_PATH)


def safe_fetch_session_responses(session_code: str) -> list[AssessmentResponse]:
    try:
        return fetch_session_responses(DB_PATH, session_code)
    except sqlite3.OperationalError:
        # Handles upgraded deployments where the process restarted before schema migration.
        ensure_assessment_schema()
        return fetch_session_responses(DB_PATH, session_code)


def safe_store_assessment_response(response: AssessmentResponse) -> None:
    try:
        store_assessment_response(DB_PATH, response)
    except sqlite3.OperationalError:
        ensure_assessment_schema()
        store_assessment_response(DB_PATH, response)


def correct_answers_by_question_for_session(session_code: str, session: RuntimeSession) -> dict[str, str]:
    answers: dict[str, str] = {}

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT id, correct_option_id
        FROM saved_quizzes
        WHERE session_code = ?
        """,
        (session_code,),
    )
    for row in cursor.fetchall():
        answers[f"quiz-{row['id']}"] = str(row["correct_option_id"])
    conn.close()

    if session.current_quiz is not None:
        answers[current_assessment_question_id(session)] = session.current_quiz.correct_option_id

    return answers


def generate_quizes_analyticss_pdf(session_code: str, session: RuntimeSession) -> bytes:
    if canvas is None or A4 is None or cm is None:
        raise RuntimeError("PDF generation dependency is not available")

    responses = safe_fetch_session_responses(session_code)
    metrics = calculate_session_metrics(responses)
    breakdown = question_performance_breakdown(responses)
    correct_by_question = correct_answers_by_question_for_session(session_code, session)
    profiles = build_individual_profiles(responses, correct_by_question)

    with TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        chart_path: Path | None = None
        if breakdown:
            chart_path = temp_path / "class-question-performance.png"
            save_question_performance_bar_chart(
                breakdown,
                chart_path,
                title="Quizes Analyticss - Class Performance by Question",
            )

        student_table_paths: list[tuple[str, Path]] = []
        for student_id, student_rows in profiles.items():
            if not student_rows:
                continue
            table_path = temp_path / f"{student_id}-summary.png"
            save_student_summary_table(student_id, student_rows, table_path)
            student_table_paths.append((student_id, table_path))

        buffer = io.BytesIO()
        pdf = canvas.Canvas(buffer, pagesize=A4)
        page_width, page_height = A4
        margin_x = 2 * cm
        y = page_height - 2 * cm

        pdf.setFont("Helvetica-Bold", 22)
        pdf.drawString(margin_x, y, "Quizes Analyticss")
        y -= 22
        pdf.setFont("Helvetica", 10)
        pdf.drawString(margin_x, y, f"Session: {session_code}")
        y -= 14
        pdf.drawString(margin_x, y, f"Generated at: {now_iso()}")

        y -= 22
        pdf.setFont("Helvetica-Bold", 13)
        pdf.drawString(margin_x, y, "Session Metrics")
        y -= 16
        pdf.setFont("Helvetica", 10)
        pdf.drawString(margin_x, y, f"Total responses: {metrics['total_responses']}")
        y -= 14
        pdf.drawString(margin_x, y, f"Total accuracy: {metrics['total_accuracy_percentage']}%")
        y -= 14
        pdf.drawString(
            margin_x,
            y,
            f"Average time per question: {metrics['average_time_per_question_seconds']}s",
        )
        y -= 14
        pdf.drawString(
            margin_x,
            y,
            (
                "Answer frequency: "
                f"correct={metrics['answer_frequency']['correct']}, "
                f"incorrect={metrics['answer_frequency']['incorrect']}"
            ),
        )

        if chart_path and chart_path.exists() and ImageReader is not None:
            y -= 22
            pdf.setFont("Helvetica-Bold", 13)
            pdf.drawString(margin_x, y, "Class Performance by Question")
            y -= 12
            image_top = y
            image_height = 10 * cm
            image_width = page_width - (2 * margin_x)
            pdf.drawImage(
                ImageReader(str(chart_path)),
                margin_x,
                image_top - image_height,
                width=image_width,
                height=image_height,
                preserveAspectRatio=True,
                anchor="n",
            )

        for student_id, table_path in student_table_paths:
            pdf.showPage()
            pdf.setFont("Helvetica-Bold", 16)
            pdf.drawString(margin_x, page_height - 2 * cm, f"Student Summary - {student_id}")
            if table_path.exists() and ImageReader is not None:
                image_width = page_width - (2 * margin_x)
                image_height = page_height - (4 * cm)
                pdf.drawImage(
                    ImageReader(str(table_path)),
                    margin_x,
                    2 * cm,
                    width=image_width,
                    height=image_height,
                    preserveAspectRatio=True,
                    anchor="sw",
                )

        if not responses:
            y -= 24
            pdf.setFont("Helvetica", 11)
            pdf.drawString(margin_x, y, "No quiz responses were captured for this session.")

        pdf.save()
        return buffer.getvalue()


def generate_quiz_information_xlsx(session_code: str, session: RuntimeSession) -> bytes:
    responses = safe_fetch_session_responses(session_code)
    metrics = calculate_session_metrics(responses)
    breakdown = question_performance_breakdown(responses)
    correct_by_question = correct_answers_by_question_for_session(session_code, session)
    profiles = build_individual_profiles(responses, correct_by_question)

    quiz_bank_rows: list[list[Any]] = []
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT id, session_code, question, options_json, correct_option_id, created_at
        FROM saved_quizzes
        WHERE session_code = ?
        ORDER BY datetime(created_at) ASC
        """,
        (session_code,),
    )
    for row in cursor.fetchall():
        options_raw = json.loads(row["options_json"])
        options = {str(item.get("id", "")).upper(): str(item.get("text", "")) for item in options_raw}
        quiz_bank_rows.append([
            row["id"],
            row["session_code"],
            row["question"],
            options.get("A", ""),
            options.get("B", ""),
            options.get("C", ""),
            options.get("D", ""),
            row["correct_option_id"],
            row["created_at"],
        ])
    conn.close()

    try:
        from openpyxl import Workbook

        workbook = Workbook()
        metrics_sheet = workbook.active
        metrics_sheet.title = "Session Metrics"
        metrics_sheet.append(["Metric", "Value"])
        metrics_sheet.append(["session_code", session_code])
        metrics_sheet.append(["generated_at", now_iso()])
        metrics_sheet.append(["total_responses", metrics["total_responses"]])
        metrics_sheet.append(["total_accuracy_percentage", metrics["total_accuracy_percentage"]])
        metrics_sheet.append([
            "average_time_per_question_seconds",
            metrics["average_time_per_question_seconds"],
        ])
        metrics_sheet.append(["correct_answers", metrics["answer_frequency"]["correct"]])
        metrics_sheet.append(["incorrect_answers", metrics["answer_frequency"]["incorrect"]])

        performance_sheet = workbook.create_sheet("Question Performance")
        performance_sheet.append([
            "question_id",
            "total_answers",
            "correct_answers",
            "incorrect_answers",
            "accuracy_percentage",
        ])
        for breakdown_row in breakdown:
            performance_sheet.append([
                breakdown_row["question_id"],
                breakdown_row["total_answers"],
                breakdown_row["correct_answers"],
                breakdown_row["incorrect_answers"],
                breakdown_row["accuracy_percentage"],
            ])

        responses_sheet = workbook.create_sheet("Student Responses")
        responses_sheet.append([
            "student_id",
            "question_id",
            "timestamp",
            "raw_answer",
            "is_correct",
        ])
        for response_row in responses:
            responses_sheet.append([
                response_row.student_id,
                response_row.question_id,
                response_row.timestamp,
                response_row.raw_answer,
                "TRUE" if response_row.is_correct else "FALSE",
            ])

        profiles_sheet = workbook.create_sheet("Student Profiles")
        profiles_sheet.append([
            "student_id",
            "question_id",
            "student_answer",
            "correct_answer",
            "is_correct",
            "timestamp",
        ])
        for student_id, student_rows in profiles.items():
            for profile_row in student_rows:
                profiles_sheet.append([
                    student_id,
                    profile_row["question_id"],
                    profile_row["student_answer"],
                    profile_row["correct_answer"],
                    "TRUE" if profile_row["is_correct"] else "FALSE",
                    profile_row["timestamp"],
                ])

        quiz_sheet = workbook.create_sheet("Quiz Bank")
        quiz_sheet.append([
            "quiz_id",
            "session_code",
            "question",
            "option_a",
            "option_b",
            "option_c",
            "option_d",
            "correct_option_id",
            "created_at",
        ])
        for quiz_row in quiz_bank_rows:
            quiz_sheet.append(quiz_row)

        output = BytesIO()
        workbook.save(output)
        return output.getvalue()
    except Exception:
        # Fallback when openpyxl is unavailable in runtime environment.
        try:
            import xlsxwriter
        except Exception as exc:
            raise RuntimeError("XLSX generation is unavailable. Install openpyxl or xlsxwriter.") from exc

        output = BytesIO()
        workbook = xlsxwriter.Workbook(output, {"in_memory": True})

        metrics_sheet = workbook.add_worksheet("Session Metrics")
        metrics_rows = [
            ["Metric", "Value"],
            ["session_code", session_code],
            ["generated_at", now_iso()],
            ["total_responses", metrics["total_responses"]],
            ["total_accuracy_percentage", metrics["total_accuracy_percentage"]],
            ["average_time_per_question_seconds", metrics["average_time_per_question_seconds"]],
            ["correct_answers", metrics["answer_frequency"]["correct"]],
            ["incorrect_answers", metrics["answer_frequency"]["incorrect"]],
        ]
        for row_index, row_values in enumerate(metrics_rows):
            metrics_sheet.write_row(row_index, 0, row_values)

        performance_sheet = workbook.add_worksheet("Question Performance")
        performance_sheet.write_row(0, 0, [
            "question_id",
            "total_answers",
            "correct_answers",
            "incorrect_answers",
            "accuracy_percentage",
        ])
        for idx, breakdown_row in enumerate(breakdown, start=1):
            performance_sheet.write_row(idx, 0, [
                breakdown_row["question_id"],
                breakdown_row["total_answers"],
                breakdown_row["correct_answers"],
                breakdown_row["incorrect_answers"],
                breakdown_row["accuracy_percentage"],
            ])

        responses_sheet = workbook.add_worksheet("Student Responses")
        responses_sheet.write_row(0, 0, [
            "student_id",
            "question_id",
            "timestamp",
            "raw_answer",
            "is_correct",
        ])
        for idx, response_row in enumerate(responses, start=1):
            responses_sheet.write_row(idx, 0, [
                response_row.student_id,
                response_row.question_id,
                response_row.timestamp,
                response_row.raw_answer,
                "TRUE" if response_row.is_correct else "FALSE",
            ])

        profiles_sheet = workbook.add_worksheet("Student Profiles")
        profiles_sheet.write_row(0, 0, [
            "student_id",
            "question_id",
            "student_answer",
            "correct_answer",
            "is_correct",
            "timestamp",
        ])
        profile_line = 1
        for student_id, student_rows in profiles.items():
            for profile_row in student_rows:
                profiles_sheet.write_row(profile_line, 0, [
                    student_id,
                    profile_row["question_id"],
                    profile_row["student_answer"],
                    profile_row["correct_answer"],
                    "TRUE" if profile_row["is_correct"] else "FALSE",
                    profile_row["timestamp"],
                ])
                profile_line += 1

        quiz_sheet = workbook.add_worksheet("Quiz Bank")
        quiz_sheet.write_row(0, 0, [
            "quiz_id",
            "session_code",
            "question",
            "option_a",
            "option_b",
            "option_c",
            "option_d",
            "correct_option_id",
            "created_at",
        ])
        for idx, quiz_row in enumerate(quiz_bank_rows, start=1):
            quiz_sheet.write_row(idx, 0, quiz_row)

        workbook.close()
        return output.getvalue()


def generate_report_bundle_zip(code: str, session: RuntimeSession) -> bytes:
    report = session.final_report if session.final_report is not None else build_full_analytics_report(session)
    if session.final_report_insights is None:
        session.final_report_insights = build_analytics_insights_with_ai(report)

    legacy_pdf = generate_session_report_pdf(report, session.final_report_insights)
    quizes_analyticss_pdf = generate_quizes_analyticss_pdf(code, session)
    quizzes_xlsx = generate_quiz_information_xlsx(code, session)

    archive = BytesIO()
    with zipfile.ZipFile(archive, mode="w", compression=zipfile.ZIP_DEFLATED) as zip_file:
        zip_file.writestr(f"session-{code}-analytics-report.pdf", legacy_pdf)
        zip_file.writestr(f"session-{code}-Quizes Analyticss.pdf", quizes_analyticss_pdf)
        zip_file.writestr(f"session-{code}-quizes-information.xlsx", quizzes_xlsx)

    return archive.getvalue()


app = FastAPI(title="VIA Live", version="1.0.0")


def parse_allowed_origins(raw: str) -> list[str]:
    origins: list[str] = []
    for origin in raw.split(","):
        normalized = origin.strip().rstrip("/")
        if normalized:
            origins.append(normalized)
    return origins


allowed_origins = parse_allowed_origins(
    os.getenv("ALLOWED_ORIGINS", "http://localhost:5173,http://127.0.0.1:5173")
)
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_methods=["*"],
    allow_headers=["*"],
    allow_credentials=True,
)


@app.on_event("startup")
def startup() -> None:
    init_db()


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok", "time": now_iso()}


@app.post("/api/auth/register", response_model=AuthResponse)
def register(payload: AuthRegisterRequest) -> AuthResponse:
    email = payload.email.strip().lower()
    display_name = payload.display_name.strip()
    password = payload.password
    role = (payload.role or "teacher").strip().lower()

    if role not in {"teacher", "student"}:
        raise HTTPException(status_code=400, detail="role must be teacher or student")
    if not email or "@" not in email:
        raise HTTPException(status_code=400, detail="A valid email is required")
    if len(password) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")
    if not display_name:
        raise HTTPException(status_code=400, detail="display_name is required")

    salt_hex, password_hash_hex = hash_password(password)
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    try:
        cursor.execute(
            """
            INSERT INTO users(email, display_name, role, password_salt, password_hash, created_at)
            VALUES (?, ?, ?, ?, ?, ?)
            """,
            (email, display_name, role, salt_hex, password_hash_hex, now_iso()),
        )
    except sqlite3.IntegrityError:
        conn.close()
        raise HTTPException(status_code=409, detail="Email is already registered")

    user_id = cursor.lastrowid
    token = create_auth_token()
    cursor.execute(
        "INSERT INTO auth_tokens(token, user_id, created_at) VALUES (?, ?, ?)",
        (token, user_id, now_iso()),
    )
    conn.commit()

    cursor.execute("SELECT id, email, display_name, role FROM users WHERE id = ?", (user_id,))
    row = cursor.fetchone()
    conn.close()

    if not row:
        raise HTTPException(status_code=500, detail="User creation failed")
    user = UserPublic(**dict(row))
    return AuthResponse(token=token, user=user)


@app.post("/api/auth/login", response_model=AuthResponse)
def login(payload: AuthLoginRequest) -> AuthResponse:
    email = payload.email.strip().lower()
    password = payload.password

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        "SELECT id, email, display_name, role, password_salt, password_hash FROM users WHERE email = ?",
        (email,),
    )
    row = cursor.fetchone()
    if not row:
        conn.close()
        raise HTTPException(status_code=401, detail="Invalid email or password")

    salt_hex, expected_hash_hex = row["password_salt"], row["password_hash"]
    _, actual_hash_hex = hash_password(password, salt_hex)
    if actual_hash_hex != expected_hash_hex:
        conn.close()
        raise HTTPException(status_code=401, detail="Invalid email or password")

    token = create_auth_token()
    cursor.execute(
        "INSERT INTO auth_tokens(token, user_id, created_at) VALUES (?, ?, ?)",
        (token, row["id"], now_iso()),
    )
    conn.commit()
    conn.close()

    user = UserPublic(
        id=row["id"],
        email=row["email"],
        display_name=row["display_name"],
        role=row["role"],
    )
    return AuthResponse(token=token, user=user)


@app.get("/api/auth/me", response_model=UserPublic)
def auth_me(authorization: str | None = Header(default=None)) -> UserPublic:
    user = require_user(authorization)
    return UserPublic(**user)


# ---------------------------------------------------------------------------
# GitHub OAuth
# ---------------------------------------------------------------------------

BACKEND_URL = os.getenv("BACKEND_URL", "http://localhost:9000").strip().rstrip("/")


@app.get("/api/auth/oauth/github")
def github_oauth_start() -> RedirectResponse:
    if not GITHUB_CLIENT_ID:
        raise HTTPException(status_code=503, detail="GitHub OAuth is not configured")
    params = urlencode({
        "client_id": GITHUB_CLIENT_ID,
        "redirect_uri": f"{BACKEND_URL}/api/auth/oauth/github/callback",
        "scope": "user:email",
    })
    return RedirectResponse(f"https://github.com/login/oauth/authorize?{params}")


@app.get("/api/auth/oauth/github/callback")
def github_oauth_callback(code: str | None = None, error: str | None = None) -> RedirectResponse:
    if error or not code:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=github_denied")
    if not GITHUB_CLIENT_ID or not GITHUB_CLIENT_SECRET:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=not_configured")

    # Exchange code for access token
    try:
        token_req = Request(
            "https://github.com/login/oauth/access_token",
            data=json.dumps({
                "client_id": GITHUB_CLIENT_ID,
                "client_secret": GITHUB_CLIENT_SECRET,
                "code": code,
            }).encode(),
            headers={"Accept": "application/json", "Content-Type": "application/json"},
            method="POST",
        )
        with urlopen(token_req, timeout=10) as resp:
            token_data = json.loads(resp.read())
    except Exception:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=github_token_failed")

    access_token = token_data.get("access_token")
    if not access_token:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=github_no_token")

    # Get user info
    try:
        user_req = Request(
            "https://api.github.com/user",
            headers={"Authorization": f"Bearer {access_token}", "Accept": "application/vnd.github+json"},
        )
        with urlopen(user_req, timeout=10) as resp:
            gh_user = json.loads(resp.read())
    except Exception:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=github_user_failed")

    # Get primary verified email if not public
    email = gh_user.get("email")
    if not email:
        try:
            email_req = Request(
                "https://api.github.com/user/emails",
                headers={"Authorization": f"Bearer {access_token}", "Accept": "application/vnd.github+json"},
            )
            with urlopen(email_req, timeout=10) as resp:
                emails = json.loads(resp.read())
            for e in emails:
                if e.get("primary") and e.get("verified"):
                    email = e["email"]
                    break
        except Exception:
            pass

    if not email:
        return RedirectResponse(f"{FRONTEND_URL}?oauth_error=github_no_email")

    display_name = gh_user.get("name") or gh_user.get("login") or email.split("@")[0]
    user, auth_token = upsert_oauth_user(email.strip().lower(), display_name)
    return RedirectResponse(f"{FRONTEND_URL}?oauth_token={auth_token}&oauth_user={json.dumps(user)}")


# ---------------------------------------------------------------------------
# Google OAuth (access token → userinfo)
# ---------------------------------------------------------------------------

class GoogleTokenRequest(BaseModel):
    access_token: str  # Google OAuth access token from frontend implicit flow


@app.post("/api/auth/oauth/google", response_model=AuthResponse)
def google_oauth(payload: GoogleTokenRequest) -> AuthResponse:
    try:
        req = Request(
            "https://www.googleapis.com/oauth2/v3/userinfo",
            headers={"Authorization": f"Bearer {payload.access_token}"},
        )
        with urlopen(req, timeout=10) as resp:
            info = json.loads(resp.read())
    except HTTPError as exc:
        raise HTTPException(status_code=401, detail="Invalid Google token") from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail="Google token verification failed") from exc

    email = info.get("email")
    if not email or not info.get("email_verified"):
        raise HTTPException(status_code=400, detail="Google account email not verified")

    display_name = info.get("name") or email.split("@")[0]
    user, token = upsert_oauth_user(email.strip().lower(), display_name)
    return AuthResponse(token=token, user=UserPublic(**user))


@app.get("/api/library/sessions")
def list_library_sessions(authorization: str | None = Header(default=None)) -> dict[str, Any]:
    user = require_user(authorization)
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        """
        SELECT code, created_at, teacher_name, active
        FROM sessions
        WHERE owner_user_id = ? OR (owner_user_id IS NULL AND teacher_name = ?)
        ORDER BY datetime(created_at) DESC
        LIMIT 100
        """,
        (user["id"], user["display_name"]),
    )
    rows = [dict(row) for row in cursor.fetchall()]
    conn.close()
    return {"sessions": rows}


@app.post("/api/presentations", response_model=PresentationItem)
async def upload_presentation(
    file: UploadFile = File(...),
    session_code: str = Form(""),
    authorization: str | None = Header(default=None),
) -> PresentationItem:
    user = require_user(authorization)

    normalized_session_code = session_code.strip().upper()
    if not normalized_session_code:
        raise HTTPException(status_code=400, detail="session_code is required")

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    cursor.execute(
        "SELECT code, teacher_name, owner_user_id FROM sessions WHERE code = ?",
        (normalized_session_code,),
    )
    session_row = cursor.fetchone()
    conn.close()
    if not session_row:
        raise HTTPException(status_code=404, detail="Session not found")
    session_owner_id = session_row["owner_user_id"]
    is_owner = (
        (session_owner_id is not None and session_owner_id == user["id"])
        or (session_owner_id is None and session_row["teacher_name"] == user["display_name"])
    )
    if not is_owner:
        raise HTTPException(status_code=403, detail="You can only upload files for your own sessions")

    original_name = (file.filename or "presentation").strip() or "presentation"
    extension = Path(original_name).suffix
    safe_extension = extension[:10] if extension else ""
    stored_name = f"{uuid.uuid4().hex}{safe_extension}"
    file_path = UPLOADS_DIR / stored_name

    raw = await file.read()
    size_bytes = len(raw)
    if size_bytes <= 0:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")
    if size_bytes > 50 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File exceeds 50MB limit")

    file_path.write_bytes(raw)
    created_at = now_iso()

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute(
        """
        INSERT INTO presentations(user_id, session_code, original_name, stored_name, mime_type, size_bytes, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (
            user["id"],
            normalized_session_code,
            original_name,
            stored_name,
            file.content_type or "application/octet-stream",
            size_bytes,
            created_at,
        ),
    )
    presentation_id = cursor.lastrowid
    conn.commit()
    conn.close()

    return PresentationItem(
        id=presentation_id,
        session_code=normalized_session_code,
        original_name=original_name,
        mime_type=file.content_type or "application/octet-stream",
        size_bytes=size_bytes,
        created_at=created_at,
        download_url=f"/api/presentations/{presentation_id}/download",
    )


@app.get("/api/presentations")
def list_presentations(
    session_code: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
) -> dict[str, Any]:
    user = require_user(authorization)
    normalized_session_code = (session_code or "").strip().upper() or None

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()

    if normalized_session_code:
        cursor.execute(
            "SELECT teacher_name, owner_user_id FROM sessions WHERE code = ?",
            (normalized_session_code,),
        )
        session_row = cursor.fetchone()
        if not session_row:
            conn.close()
            raise HTTPException(status_code=404, detail="Session not found")

        session_owner_id = session_row["owner_user_id"]
        is_session_owner = (
            (session_owner_id is not None and session_owner_id == user["id"])
            or (session_owner_id is None and session_row["teacher_name"] == user["display_name"])
        )
        if is_session_owner:
            cursor.execute(
                """
                                SELECT p.id, p.session_code, p.original_name, p.mime_type, p.size_bytes, p.created_at
                                FROM presentations p
                WHERE p.user_id = ?
                  AND p.session_code = ?
                ORDER BY datetime(p.created_at) DESC
                """,
                                (user["id"], normalized_session_code),
            )
        else:
            cursor.execute(
                """
                SELECT p.id, p.session_code, p.original_name, p.mime_type, p.size_bytes, p.created_at
                FROM presentations p
                JOIN sessions s ON s.code = p.session_code
                JOIN users u ON u.id = p.user_id
                WHERE p.session_code = ?
                                    AND (
                                        (s.owner_user_id IS NOT NULL AND s.owner_user_id = u.id)
                                        OR (s.owner_user_id IS NULL AND s.teacher_name = u.display_name)
                                    )
                ORDER BY datetime(p.created_at) DESC
                """,
                (normalized_session_code,),
            )
    else:
        cursor.execute(
            """
            SELECT id, session_code, original_name, mime_type, size_bytes, created_at
            FROM presentations
            WHERE user_id = ?
            ORDER BY datetime(created_at) DESC
            """,
            (user["id"],),
        )

    items = [
        PresentationItem(
            id=row["id"],
            session_code=row["session_code"],
            original_name=row["original_name"],
            mime_type=row["mime_type"] or "application/octet-stream",
            size_bytes=row["size_bytes"],
            created_at=row["created_at"],
            download_url=f"/api/presentations/{row['id']}/download",
        ).model_dump()
        for row in cursor.fetchall()
    ]
    conn.close()
    return {"presentations": items}


@app.get("/api/presentations/{presentation_id}/download")
def download_presentation(
    presentation_id: int,
    session_code: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
) -> FileResponse:
    user = require_user(authorization)
    normalized_session_code = (session_code or "").strip().upper() or None

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    row = get_accessible_presentation_row(
        cursor,
        user=user,
        presentation_id=presentation_id,
        normalized_session_code=normalized_session_code,
    )
    conn.close()

    if not row:
        raise HTTPException(status_code=404, detail="Presentation not found")

    file_path = UPLOADS_DIR / row["stored_name"]
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Stored file is missing")

    return FileResponse(
        path=file_path,
        filename=row["original_name"],
        media_type=row["mime_type"] or "application/octet-stream",
    )


@app.post("/api/presentations/{presentation_id}/notes-png")
def generate_presentation_notes_png(
    presentation_id: int,
    session_code: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
) -> Response:
    user = require_user(authorization)
    normalized_session_code = (session_code or "").strip().upper() or None

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()
    row = get_accessible_presentation_row(
        cursor,
        user=user,
        presentation_id=presentation_id,
        normalized_session_code=normalized_session_code,
    )
    conn.close()

    if not row:
        raise HTTPException(status_code=404, detail="Presentation not found")

    file_path = UPLOADS_DIR / row["stored_name"]
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Stored file is missing")

    try:
        extracted_text = extract_text_from_presentation(
            file_path,
            row["original_name"],
            row["mime_type"],
        )
        notes_text = build_student_notes_with_ai(extracted_text, row["original_name"])
        png_bytes = render_notes_png(
            title=f"Student Notes: {Path(row['original_name']).stem}",
            notes_text=notes_text,
        )
    except RuntimeError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    safe_stem = Path(row["original_name"]).stem.replace('"', "").strip()[:70] or "presentation"
    download_name = f"{safe_stem}-student-notes.png"
    return Response(
        content=png_bytes,
        media_type="image/png",
        headers={
            "Content-Disposition": f'attachment; filename="{download_name}"',
            "Cache-Control": "no-store",
        },
    )


def save_quiz_for_user(
    *,
    user_id: int,
    session_code: str | None,
    quiz: QuizPayload,
) -> SavedQuizItem:
    question = quiz.question.strip()
    if not question:
        raise HTTPException(status_code=400, detail="question is required")

    correct_option_id = quiz.correct_option_id.strip().upper()
    options = quiz.options
    valid_option_ids = {option.id.upper() for option in options}
    if correct_option_id not in valid_option_ids:
        raise HTTPException(status_code=400, detail="correct_option_id must match an option id")

    serialized_options = json.dumps([option.model_dump() for option in options])
    normalized_session_code = (session_code or "").strip().upper() or None

    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    created_at = now_iso()
    cursor.execute(
        """
        INSERT INTO saved_quizzes(user_id, session_code, question, options_json, correct_option_id, created_at)
        VALUES (?, ?, ?, ?, ?, ?)
        """,
        (user_id, normalized_session_code, question, serialized_options, correct_option_id, created_at),
    )
    quiz_id = cursor.lastrowid
    conn.commit()
    conn.close()

    return SavedQuizItem(
        id=quiz_id,
        session_code=normalized_session_code,
        question=question,
        options=options,
        correct_option_id=correct_option_id,
        created_at=created_at,
    )


@app.post("/api/quizzes/save", response_model=SavedQuizItem)
def save_quiz(payload: SavedQuizCreateRequest, authorization: str | None = Header(default=None)) -> SavedQuizItem:
    user = require_user(authorization)
    quiz = QuizPayload(
        question=payload.question,
        options=payload.options,
        correct_option_id=payload.correct_option_id,
    )
    return save_quiz_for_user(
        user_id=user["id"],
        session_code=payload.session_code,
        quiz=quiz,
    )


@app.get("/api/quizzes")
def list_saved_quizzes(
    session_code: str | None = Query(default=None),
    authorization: str | None = Header(default=None),
) -> dict[str, Any]:
    user = require_user(authorization)
    normalized_session_code = (session_code or "").strip().upper() or None

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cursor = conn.cursor()

    if normalized_session_code:
        cursor.execute(
            "SELECT teacher_name, owner_user_id FROM sessions WHERE code = ?",
            (normalized_session_code,),
        )
        session_row = cursor.fetchone()
        if not session_row:
            conn.close()
            raise HTTPException(status_code=404, detail="Session not found")

        session_owner_id = session_row["owner_user_id"]
        is_session_owner = (
            (session_owner_id is not None and session_owner_id == user["id"])
            or (session_owner_id is None and session_row["teacher_name"] == user["display_name"])
        )

        if is_session_owner:
            cursor.execute(
                """
                SELECT id, session_code, question, options_json, correct_option_id, created_at
                FROM saved_quizzes
                WHERE user_id = ? AND session_code = ?
                ORDER BY datetime(created_at) DESC
                """,
                (user["id"], normalized_session_code),
            )
        else:
            cursor.execute(
                """
                SELECT q.id, q.session_code, q.question, q.options_json, q.correct_option_id, q.created_at
                FROM saved_quizzes q
                JOIN sessions s ON s.code = q.session_code
                JOIN users u ON u.id = q.user_id
                WHERE q.session_code = ?
                  AND (
                    (s.owner_user_id IS NOT NULL AND s.owner_user_id = u.id)
                    OR (s.owner_user_id IS NULL AND s.teacher_name = u.display_name)
                  )
                ORDER BY datetime(q.created_at) DESC
                """,
                (normalized_session_code,),
            )
    else:
        cursor.execute(
            """
            SELECT id, session_code, question, options_json, correct_option_id, created_at
            FROM saved_quizzes
            WHERE user_id = ?
            ORDER BY datetime(created_at) DESC
            """,
            (user["id"],),
        )

    rows = cursor.fetchall()
    conn.close()

    live_saved_quiz_ids = {
        live_session.current_quiz_saved_id
        for live_session in SESSIONS.values()
        if live_session.current_quiz_saved_id is not None
        and not live_session.quiz_voting_closed
        and not live_session.quiz_hidden
    }

    quizzes: list[dict[str, Any]] = []
    for row in rows:
        options_raw = json.loads(row["options_json"])
        is_live = row["id"] in live_saved_quiz_ids
        answer_revealed = not is_live
        quizzes.append(
            SavedQuizListItem(
                id=row["id"],
                session_code=row["session_code"],
                question=row["question"],
                options=[QuizOption(**item) for item in options_raw],
                correct_option_id=row["correct_option_id"] if answer_revealed else None,
                answer_revealed=answer_revealed,
                is_live=is_live,
                created_at=row["created_at"],
            ).model_dump()
        )

    return {"quizzes": quizzes}


@app.post("/api/sessions", response_model=SessionCreateResponse)
def create_session(payload: SessionCreateRequest, authorization: str | None = Header(default=None)) -> SessionCreateResponse:
    teacher_name = payload.teacher_name.strip()
    if not teacher_name:
        raise HTTPException(status_code=400, detail="teacher_name is required")

    user = get_user_by_token(parse_bearer_token(authorization))
    owner_user_id = user["id"] if user else None

    code = generate_session_code()
    SESSIONS[code] = RuntimeSession(code=code, teacher_name=teacher_name)
    record_engagement_point(SESSIONS[code], "session_created")
    insert_session(code, teacher_name, owner_user_id=owner_user_id)
    return SessionCreateResponse(code=code)


@app.get("/api/sessions/{code}/analytics")
def session_analytics(code: str) -> dict[str, Any]:
    session = SESSIONS.get(code.upper())
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    return analytics_for_session(session)


@app.post("/api/sessions/{code}/end")
async def end_session(code: str) -> dict[str, Any]:
    code = code.upper()
    session = SESSIONS.get(code)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")

    if session.final_report is not None:
        return session.final_report

    session.active = False
    session.ended_at_epoch = time.time()
    record_engagement_point(session, "session_ended")
    set_session_active(code, False)

    report = build_full_analytics_report(session)
    session.final_report = report

    insert_event(
        code,
        "session_end",
        {
            "generated_at": report["generated_at"],
            "engagement_score": report["analytics"]["engagement"]["score"],
        },
    )

    await broadcast(
        session,
        {
            "type": "session_ended",
            "payload": report,
        },
    )

    for client in list(session.clients.values()):
        try:
            await client.websocket.close(code=1000, reason="Session ended by teacher")
        except Exception:
            pass
    session.clients.clear()

    return report


@app.get("/api/sessions/{code}/report.pdf")
def download_session_report_pdf(code: str) -> Response:
    code = code.upper()
    session = SESSIONS.get(code)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")

    report = session.final_report if session.final_report is not None else build_full_analytics_report(session)
    if session.final_report_insights is None:
        session.final_report_insights = build_analytics_insights_with_ai(report)

    pdf_bytes = generate_session_report_pdf(report, session.final_report_insights)
    filename = f"session-{code}-analytics-report.pdf"
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.get("/api/sessions/{code}/report-bundle.zip")
def download_session_report_bundle(code: str) -> Response:
    code = code.upper()
    session = SESSIONS.get(code)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")

    bundle_bytes = generate_report_bundle_zip(code, session)
    filename = f"session-{code}-reports-bundle.zip"
    return Response(
        content=bundle_bytes,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@app.websocket("/ws/{code}")
async def websocket_room(websocket: WebSocket, code: str, role: str, name: str) -> None:
    code = code.upper()
    role = role.lower().strip()
    name = name.strip()

    if role not in {"teacher", "student"} or not name:
        await websocket.close(code=1008, reason="Invalid role or name")
        return

    session = SESSIONS.get(code)
    if not session or not session.active:
        await websocket.close(code=1008, reason="Unknown session code")
        return

    await websocket.accept()

    if role == "teacher" and get_teacher(session):
        await websocket.close(code=1008, reason="Teacher already connected")
        return

    client_id = str(uuid.uuid4())
    state = ClientState(client_id=client_id, websocket=websocket, role=role, name=name)
    session.clients[client_id] = state
    record_engagement_point(session, "participant_joined")

    insert_event(code, "join", {"client_id": client_id, "role": role, "name": name})

    await send_json(
        websocket,
        {
            "type": "welcome",
            "payload": {
                "client_id": client_id,
                "session_code": code,
                **session_state_payload(session),
            },
        },
    )

    await send_json(
        websocket,
        {
            "type": "session_state",
            "payload": session_state_payload(session),
        },
    )

    await broadcast(
        session,
        {
            "type": "participant_joined",
            "payload": {"client_id": client_id, "role": role, "name": name},
        },
    )

    await broadcast(
        session,
        {
            "type": "metrics",
            "payload": metrics_payload(session),
        },
    )

    if role == "student":
        teacher = get_teacher(session)
        if teacher:
            await send_json(
                teacher.websocket,
                {
                    "type": "student_joined",
                    "payload": {"student_id": client_id, "name": name},
                },
            )
    elif role == "teacher":
        await send_json(
            websocket,
            {
                "type": "anonymous_questions",
                "payload": anonymous_questions_payload(session),
            },
        )

    try:
        while True:
            raw = await websocket.receive_text()
            message = json.loads(raw)
            msg_type = message.get("type")
            payload = message.get("payload", {})

            # Detect natural break expiry and reset the focus-period timer
            if session.break_active_until and time.time() > session.break_active_until:
                session.break_active_until = None
                session.break_votes.clear()
                session.focus_period_ends_at = time.time() + FOCUS_PERIOD_SECONDS
                await broadcast(session, {"type": "break_ended", "payload": {"reason": "expired"}})
                await broadcast(
                    session,
                    {
                        "type": "focus_timer_reset",
                        "payload": {"focus_period_ends_at": session.focus_period_ends_at},
                    },
                )
                await broadcast(session, {"type": "metrics", "payload": metrics_payload(session)})

            if msg_type == "signal":
                target_id = payload.get("target_id")
                target = session.clients.get(target_id)
                if target:
                    await send_json(
                        target.websocket,
                        {
                            "type": "signal",
                            "payload": {
                                "from_id": client_id,
                                "description": payload.get("description"),
                                "candidate": payload.get("candidate"),
                            },
                        },
                    )

            elif msg_type == "confusion":
                if role != "student":
                    continue
                try:
                    now_epoch = time.time()
                    state.last_confusion_vote_at = now_epoch
                    state.confusion_signals_sent += 1
                    state.last_active_at = now_epoch
                    insert_event(code, "confusion", {"client_id": client_id, "level": 1.0})
                    record_engagement_point(session, "confusion")
                    await broadcast(
                        session,
                        {
                            "type": "metrics",
                            "payload": metrics_payload(session),
                        },
                    )
                except Exception as e:
                    print("CONFUSION ERROR:", repr(e))

            elif msg_type == "break_vote":
                if role != "student":
                    continue

                now = time.time()

                # Reject if the 30-minute focus period hasn't elapsed yet
                if now < session.focus_period_ends_at:
                    remaining = int(session.focus_period_ends_at - now)
                    mins, secs = divmod(remaining, 60)
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {
                                "message": f"Break voting unlocks in {mins}m {secs:02d}s."
                            },
                        },
                    )
                    continue

                # Reject if a break is already active
                if session.break_active_until and now < session.break_active_until:
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {"message": "A break is already in progress."},
                        },
                    )
                    continue

                if now - state.last_break_vote_at < BREAK_COOLDOWN_SECONDS:
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {
                                "message": f"Break vote cooldown active ({BREAK_COOLDOWN_SECONDS}s)."
                            },
                        },
                    )
                    continue

                state.last_break_vote_at = now
                state.break_votes_cast += 1
                state.last_active_at = now
                session.break_votes.add(client_id)
                insert_event(code, "break_vote", {"client_id": client_id})
                record_engagement_point(session, "break_vote")

                student_count = max(session.student_count, 1)
                ratio = len(session.break_votes) / student_count

                await broadcast(
                    session,
                    {
                        "type": "metrics",
                        "payload": metrics_payload(session),
                    },
                )

                if ratio >= BREAK_THRESHOLD_PERCENT:
                    teacher = get_teacher(session)
                    if teacher:
                        await send_json(
                            teacher.websocket,
                            {
                                "type": "break_threshold_reached",
                                "payload": {"ratio": round(ratio, 3), "votes": len(session.break_votes)},
                            },
                        )

            elif msg_type == "start_break":
                if role != "teacher":
                    continue

                duration = int(payload.get("duration_seconds", 300))
                duration = max(30, min(duration, 1800))
                session.break_active_until = time.time() + duration
                session.break_votes.clear()
                # Clear focus period while break is active; it resets after break ends
                session.focus_period_ends_at = 0.0
                insert_event(code, "break_start", {"duration_seconds": duration})
                record_engagement_point(session, "break_start")

                await broadcast(
                    session,
                    {
                        "type": "break_started",
                        "payload": {
                            "end_time_epoch": session.break_active_until,
                            "focus_period_ends_at": session.focus_period_ends_at,
                        },
                    },
                )

            elif msg_type == "break_control":
                if role != "teacher":
                    continue

                action = str(payload.get("action", "")).strip().lower()

                if action == "cancel":
                    if not session.break_active_until:
                        continue
                    session.break_active_until = None
                    session.break_votes.clear()
                    # Start fresh 30-minute focus period after the break ends
                    session.focus_period_ends_at = time.time() + FOCUS_PERIOD_SECONDS
                    insert_event(code, "break_cancelled", {"by": client_id})
                    record_engagement_point(session, "break_cancelled")
                    await broadcast(session, {"type": "break_ended", "payload": {"reason": "cancelled"}})
                    await broadcast(
                        session,
                        {
                            "type": "focus_timer_reset",
                            "payload": {"focus_period_ends_at": session.focus_period_ends_at},
                        },
                    )
                    await broadcast(
                        session,
                        {
                            "type": "metrics",
                            "payload": metrics_payload(session),
                        },
                    )
                    continue

                if action == "adjust":
                    if not session.break_active_until:
                        continue

                    try:
                        delta_seconds = int(payload.get("delta_seconds", 0))
                    except Exception:
                        delta_seconds = 0

                    if delta_seconds == 0:
                        continue

                    now_epoch = time.time()
                    current_end = max(session.break_active_until, now_epoch)
                    updated_end = max(now_epoch, current_end + delta_seconds)
                    updated_end = min(updated_end, now_epoch + MAX_BREAK_DURATION_SECONDS)

                    if updated_end <= now_epoch + 1:
                        session.break_active_until = None
                        session.break_votes.clear()
                        # Start fresh 30-minute focus period after the break ends
                        session.focus_period_ends_at = now_epoch + FOCUS_PERIOD_SECONDS
                        insert_event(code, "break_adjusted", {"delta_seconds": delta_seconds, "ended": True})
                        record_engagement_point(session, "break_adjusted_end")
                        await broadcast(session, {"type": "break_ended", "payload": {"reason": "adjusted_to_zero"}})
                        await broadcast(
                            session,
                            {
                                "type": "focus_timer_reset",
                                "payload": {"focus_period_ends_at": session.focus_period_ends_at},
                            },
                        )
                        await broadcast(
                            session,
                            {
                                "type": "metrics",
                                "payload": metrics_payload(session),
                            },
                        )
                    else:
                        session.break_active_until = updated_end
                        insert_event(
                            code,
                            "break_adjusted",
                            {
                                "delta_seconds": delta_seconds,
                                "end_time_epoch": session.break_active_until,
                            },
                        )
                        record_engagement_point(session, "break_adjusted")
                        await broadcast(
                            session,
                            {
                                "type": "break_started",
                                "payload": {"end_time_epoch": session.break_active_until},
                            },
                        )
                    continue

            elif msg_type == "note_update":
                if role != "teacher":
                    continue
                session.notes = str(payload.get("text", ""))
                insert_event(code, "note_update", {"len": len(session.notes)})
                await broadcast(session, {"type": "notes", "payload": {"text": session.notes}})

            elif msg_type == "screen_share_stopped":
                if role != "teacher":
                    continue
                insert_event(code, "screen_share_stopped", {"client_id": client_id})
                await broadcast(
                    session,
                    {
                        "type": "screen_share_stopped",
                        "payload": {},
                    },
                    role="student",
                )

            elif msg_type == "ask_question":
                if role != "student":
                    continue

                text = str(payload.get("text", "")).strip()
                if not text:
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {
                                "message": "Question text is required.",
                            },
                        },
                    )
                    continue

                question = AnonymousQuestion(
                    id=uuid.uuid4().hex[:12],
                    text=text[:600],
                    created_at=now_iso(),
                )
                session.anonymous_questions.append(question)
                state.last_active_at = time.time()
                insert_event(
                    code,
                    "anonymous_question_submitted",
                    {
                        "question_id": question.id,
                        "text_len": len(question.text),
                    },
                )

                await send_json(
                    websocket,
                    {
                        "type": "anonymous_question_submitted",
                        "payload": {
                            "question_id": question.id,
                        },
                    },
                )

                teacher = get_teacher(session)
                if teacher:
                    await send_json(
                        teacher.websocket,
                        {
                            "type": "anonymous_questions",
                            "payload": anonymous_questions_payload(session),
                        },
                    )

            elif msg_type == "resolve_question":
                if role != "teacher":
                    continue

                question_id = str(payload.get("question_id", "")).strip()
                if not question_id:
                    continue

                changed = False
                for question in session.anonymous_questions:
                    if question.id != question_id:
                        continue
                    if not question.resolved:
                        question.resolved = True
                        question.resolved_at = now_iso()
                        changed = True
                    break

                if not changed:
                    continue

                insert_event(
                    code,
                    "anonymous_question_resolved",
                    {
                        "question_id": question_id,
                    },
                )
                await send_json(
                    websocket,
                    {
                        "type": "anonymous_questions",
                        "payload": anonymous_questions_payload(session),
                    },
                )

            elif msg_type == "generate_quiz":
                if role != "teacher":
                    continue
                notes_override = str(payload.get("notes", "")).strip()
                notes_input = notes_override if notes_override else session.notes
                style_preset = str(payload.get("quiz_preset", "default")).strip().lower() or "default"
                custom_prompt = str(payload.get("quiz_custom_prompt", "")).strip()
                try:
                    quiz = await asyncio.wait_for(
                        asyncio.to_thread(
                            build_quiz_with_ai,
                            notes_input,
                            style_preset,
                            custom_prompt,
                        ),
                        timeout=AI_QUIZ_GENERATION_TIMEOUT_SECONDS,
                    )
                except asyncio.TimeoutError:
                    reason = f"Quiz generation timed out after {AI_QUIZ_GENERATION_TIMEOUT_SECONDS:.0f} seconds"
                    insert_event(code, "quiz_generation_failed", {"reason": reason})
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {"message": f"Quiz generation failed: {reason}"},
                        },
                    )
                    continue
                except Exception as exc:
                    reason = str(exc)[:500]
                    insert_event(code, "quiz_generation_failed", {"reason": reason})
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {
                                "message": f"Quiz generation failed: {reason}",
                            },
                        },
                    )
                    continue
                session.current_quiz = quiz
                session.current_quiz_saved_id = None
                session.quiz_answers.clear()
                session.quiz_hidden = False
                session.quiz_cover_mode = True
                session.quiz_voting_closed = False
                session.quiz_answer_revealed = False
                conn = sqlite3.connect(DB_PATH)
                conn.row_factory = sqlite3.Row
                cursor = conn.cursor()
                cursor.execute("SELECT owner_user_id FROM sessions WHERE code = ?", (code,))
                owner_row = cursor.fetchone()
                conn.close()
                owner_user_id = owner_row["owner_user_id"] if owner_row else None
                if owner_user_id:
                    try:
                        saved_quiz = save_quiz_for_user(
                            user_id=owner_user_id,
                            session_code=code,
                            quiz=quiz,
                        )
                        session.current_quiz_saved_id = saved_quiz.id
                    except Exception:
                        # Keep live flow intact even if persistence fails.
                        session.current_quiz_saved_id = None
                record_engagement_point(session, "quiz_generated")
                insert_event(
                    code,
                    "quiz_generated",
                    {
                        **quiz.model_dump(),
                        "preset": style_preset,
                        "custom_prompt": custom_prompt,
                    },
                )
                await broadcast(session, {"type": "quiz", "payload": quiz.model_dump()})
                await broadcast(
                    session,
                    {
                        "type": "quiz_state",
                        "payload": {
                            "hidden": session.quiz_hidden,
                            "cover_mode": session.quiz_cover_mode,
                            "voting_closed": session.quiz_voting_closed,
                            "answer_revealed": session.quiz_answer_revealed,
                        },
                    },
                )

            elif msg_type == "quiz_control":
                if role != "teacher":
                    continue
                if not session.current_quiz:
                    continue

                if "hidden" in payload:
                    session.quiz_hidden = bool(payload.get("hidden"))
                if "cover_mode" in payload:
                    session.quiz_cover_mode = bool(payload.get("cover_mode"))
                if "voting_closed" in payload:
                    session.quiz_voting_closed = bool(payload.get("voting_closed"))
                if "answer_revealed" in payload:
                    session.quiz_answer_revealed = bool(payload.get("answer_revealed"))
                    # Auto-close voting when answer is revealed
                    if session.quiz_answer_revealed:
                        session.quiz_voting_closed = True

                # Compute per-option distribution for reveal broadcasts
                per_option: dict[str, dict[str, Any]] | None = None
                if session.quiz_answer_revealed and session.current_quiz:
                    total = len(session.quiz_answers)
                    per_option = {}
                    for opt in session.current_quiz.options:
                        count = sum(1 for v in session.quiz_answers.values() if v == opt.id)
                        per_option[opt.id] = {
                            "count": count,
                            "pct": round(count / total, 3) if total else 0.0,
                        }

                insert_event(
                    code,
                    "quiz_control",
                    {
                        "hidden": session.quiz_hidden,
                        "cover_mode": session.quiz_cover_mode,
                        "voting_closed": session.quiz_voting_closed,
                        "answer_revealed": session.quiz_answer_revealed,
                    },
                )
                record_engagement_point(session, "quiz_control")

                broadcast_payload: dict[str, Any] = {
                    "hidden": session.quiz_hidden,
                    "cover_mode": session.quiz_cover_mode,
                    "voting_closed": session.quiz_voting_closed,
                    "answer_revealed": session.quiz_answer_revealed,
                }
                if session.quiz_answer_revealed and session.current_quiz:
                    broadcast_payload["correct_option_id"] = session.current_quiz.correct_option_id
                    broadcast_payload["per_option"] = per_option
                await broadcast(
                    session,
                    {
                        "type": "quiz_state",
                        "payload": broadcast_payload,
                    },
                )

            elif msg_type == "quiz_answer":
                if role != "student" or not session.current_quiz or session.quiz_voting_closed:
                    continue

                option_id = str(payload.get("option_id", "")).upper()
                if option_id not in {"A", "B", "C", "D"}:
                    continue

                if client_id in session.quiz_answers:
                    continue

                session.quiz_answers[client_id] = option_id
                state.quiz_answers_submitted += 1
                if option_id == session.current_quiz.correct_option_id:
                    state.quiz_correct_answers += 1
                state.last_active_at = time.time()
                try:
                    question_id = current_assessment_question_id(session)
                    safe_store_assessment_response(
                        AssessmentResponse(
                            student_id=client_id,
                            question_id=question_id,
                            session_id=code,
                            timestamp=now_iso(),
                            raw_answer=option_id,
                            is_correct=(option_id == session.current_quiz.correct_option_id),
                        )
                    )
                except Exception:
                    # Keep the live classroom flow resilient even if analytics persistence fails.
                    pass
                insert_event(code, "quiz_answer", {"client_id": client_id, "option_id": option_id})
                record_engagement_point(session, "quiz_answer")

                summary = analytics_for_session(session)["quiz"]
                await broadcast(
                    session,
                    {
                        "type": "quiz_progress",
                        "payload": summary,
                    },
                    role="teacher",
                )

            elif msg_type == "request_analytics":
                if role != "teacher":
                    continue
                await send_json(
                    websocket,
                    {
                        "type": "analytics",
                        "payload": analytics_for_session(session),
                    },
                )

            elif msg_type == "request_state":
                await send_json(
                    websocket,
                    {
                        "type": "session_state",
                        "payload": session_state_payload(session),
                    },
                )
                if role == "teacher":
                    await send_json(
                        websocket,
                        {
                            "type": "anonymous_questions",
                            "payload": anonymous_questions_payload(session),
                        },
                    )

            elif msg_type == "explain_screen":
                if role != "student":
                    continue

                notes_override = str(payload.get("notes", "")).strip()
                notes_input = notes_override if notes_override else session.notes

                try:
                    explanation = build_screen_explanation_with_ai(notes_input)
                except Exception as exc:
                    reason = str(exc)[:500]
                    insert_event(code, "screen_explanation_failed", {"client_id": client_id, "reason": reason})
                    await send_json(
                        websocket,
                        {
                            "type": "error",
                            "payload": {
                                "message": f"Screen explanation failed: {reason}",
                            },
                        },
                    )
                    continue

                insert_event(code, "screen_explanation_generated", {"client_id": client_id})
                await send_json(
                    websocket,
                    {
                        "type": "screen_explanation",
                        "payload": {
                            "text": explanation,
                            "generated_at": now_iso(),
                        },
                    },
                )

    except WebSocketDisconnect:
        pass
    finally:
        session.clients.pop(client_id, None)
        record_engagement_point(session, "participant_left")
        insert_event(code, "leave", {"client_id": client_id, "role": role, "name": name})
        await broadcast(
            session,
            {
                "type": "participant_left",
                "payload": {"client_id": client_id, "role": role, "name": name},
            },
        )

        await broadcast(
            session,
            {
                "type": "metrics",
                "payload": metrics_payload(session),
            },
        )

        if role == "student":
            teacher = get_teacher(session)
            if teacher:
                await send_json(
                    teacher.websocket,
                    {
                        "type": "student_left",
                        "payload": {"student_id": client_id},
                    },
                )
