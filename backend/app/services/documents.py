import textwrap
from io import BytesIO
from pathlib import Path


def extract_text_from_presentation(file_path: Path, original_name: str, mime_type: str | None) -> str:
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md", ".csv", ".json"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")[:50000]

    if suffix == ".pdf" or (mime_type or "").lower() == "application/pdf":
        try:
            from pypdf import PdfReader
        except Exception as exc:
            raise RuntimeError("PDF parsing is unavailable. Install pypdf.") from exc

        reader = PdfReader(str(file_path))
        chunks: list[str] = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text)
            if sum(len(item) for item in chunks) > 50000:
                break
        return "\n\n".join(chunks)[:50000]

    if suffix == ".pptx" or "presentation" in (mime_type or "").lower():
        try:
            from pptx import Presentation
        except Exception as exc:
            raise RuntimeError("PPTX parsing is unavailable. Install python-pptx.") from exc

        prs = Presentation(str(file_path))
        chunks: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    slide_lines.append(text.strip())
            if slide_lines:
                chunks.append(f"Slide {idx}:\n" + "\n".join(slide_lines))
            if sum(len(item) for item in chunks) > 50000:
                break
        return "\n\n".join(chunks)[:50000]

    raise RuntimeError(
        f"Unsupported presentation format for AI notes: {original_name}. Supported: .pptx, .pdf, .txt, .md, .csv, .json"
    )


def render_notes_png(title: str, notes_text: str) -> bytes:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception as exc:
        raise RuntimeError("PNG rendering is unavailable. Install Pillow.") from exc

    width = 1600
    height = 2200
    margin = 90

    image = Image.new("RGB", (width, height), color=(251, 252, 255))
    draw = ImageDraw.Draw(image)

    try:
        title_font = ImageFont.truetype("arial.ttf", 54)
        body_font = ImageFont.truetype("arial.ttf", 30)
    except Exception:
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()

    current_y = margin
    draw.text((margin, current_y), title, fill=(20, 34, 62), font=title_font)
    current_y += 90
    draw.line((margin, current_y, width - margin, current_y), fill=(186, 200, 230), width=3)
    current_y += 40

    usable_width = width - (2 * margin)
    paragraphs = [line.strip() for line in notes_text.splitlines() if line.strip()]
    if not paragraphs:
        paragraphs = ["No notes could be generated for this presentation."]

    rendered_lines: list[str] = []
    for paragraph in paragraphs:
        wrapped = textwrap.wrap(paragraph, width=78) or [paragraph]
        rendered_lines.extend(wrapped)
        rendered_lines.append("")

    max_lines = 58
    if len(rendered_lines) > max_lines:
        rendered_lines = rendered_lines[: max_lines - 1] + ["... (content truncated)"]

    for line in rendered_lines:
        if line == "":
            current_y += 16
            continue

        while draw.textlength(line, font=body_font) > usable_width and len(line) > 4:
            line = line[:-2]
        draw.text((margin, current_y), line, fill=(38, 46, 66), font=body_font)
        current_y += 42
        if current_y > height - margin:
            break

    output = BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()
